"""
Universal document text extractor.

Supported formats (no extra install needed):
  PDF, DOCX, DOC, ODT, RTF, TXT, MD, HTML, XLSX, PPTX, CSV, EML

Optional (requires pytesseract + Tesseract OCR binary):
  JPG, JPEG, PNG, BMP, TIFF, GIF, WEBP
"""
from pathlib import Path


async def extract_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    fn = _EXTRACTORS.get(ext)
    if fn is None:
        raise ValueError(
            f"Format « {ext} » non supporté.\n"
            f"Formats acceptés : {', '.join(sorted(_EXTRACTORS))}"
        )
    text = fn(file_path)
    if not text or not text.strip():
        raise ValueError(f"Aucun texte extrait du fichier {Path(file_path).name}")
    return text


# ── Extractors ────────────────────────────────────────────────────────────────

def _from_pdf(path: str) -> str:
    import pdfplumber
    pages = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                pages.append(t)
    return "\n".join(pages)


def _from_docx(path: str) -> str:
    import docx as _docx
    try:
        doc = _docx.Document(path)
    except Exception:
        raise ValueError(
            "Impossible de lire ce fichier .doc/.docx.\n"
            "Astuce : ouvrez-le dans Word et enregistrez en .docx"
        )
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _from_txt(path: str) -> str:
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            with open(path, "r", encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, LookupError):
            pass
    with open(path, "rb") as f:
        return f.read().decode("utf-8", errors="replace")


def _from_rtf(path: str) -> str:
    from striprtf.striprtf import rtf_to_text
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return rtf_to_text(f.read())


def _from_html(path: str) -> str:
    from bs4 import BeautifulSoup
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    for tag in soup(["script", "style", "head", "nav", "footer"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def _from_odt(path: str) -> str:
    from odf import text as odf_text, teletype
    from odf.opendocument import load
    doc = load(path)
    texts = []
    for para in doc.getElementsByType(odf_text.P):
        t = teletype.extractText(para)
        if t.strip():
            texts.append(t)
    return "\n".join(texts)


def _from_xlsx(path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    rows = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append(" | ".join(cells))
    return "\n".join(rows)


def _from_csv(path: str) -> str:
    import csv
    rows = []
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        for row in csv.reader(f):
            cells = [c.strip() for c in row if c.strip()]
            if cells:
                rows.append(" | ".join(cells))
    return "\n".join(rows)


def _from_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
    return "\n".join(texts)


def _from_eml(path: str) -> str:
    import email
    with open(path, "rb") as f:
        msg = email.message_from_bytes(f.read())
    parts = []
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                payload = part.get_payload(decode=True)
                if payload:
                    parts.append(payload.decode("utf-8", errors="replace"))
    else:
        payload = msg.get_payload(decode=True)
        if payload:
            parts.append(payload.decode("utf-8", errors="replace"))
    return "\n".join(parts)


def _from_image(path: str) -> str:
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(path)
        # Try French + English; fallback to English only
        try:
            return pytesseract.image_to_string(img, lang="fra+eng")
        except pytesseract.TesseractError:
            return pytesseract.image_to_string(img)
    except ImportError:
        raise ValueError(
            "OCR non disponible pour les images.\n"
            "Installez : pip install pytesseract Pillow\n"
            "Et Tesseract OCR : https://tesseract-ocr.github.io/tessdoc/Installation.html"
        )


# ── Format map ────────────────────────────────────────────────────────────────

_EXTRACTORS = {
    # PDF
    ".pdf":      _from_pdf,
    # Word
    ".docx":     _from_docx,
    ".doc":      _from_docx,
    # OpenDocument
    ".odt":      _from_odt,
    # Rich text
    ".rtf":      _from_rtf,
    # Plain text
    ".txt":      _from_txt,
    ".md":       _from_txt,
    ".markdown": _from_txt,
    ".log":      _from_txt,
    ".rst":      _from_txt,
    # Web
    ".html":     _from_html,
    ".htm":      _from_html,
    # Spreadsheets
    ".xlsx":     _from_xlsx,
    ".xls":      _from_xlsx,   # openpyxl handles many .xls files
    ".csv":      _from_csv,
    ".tsv":      _from_csv,
    # Presentations
    ".pptx":     _from_pptx,
    # Email
    ".eml":      _from_eml,
    # Images (OCR — requires pytesseract + Tesseract binary)
    ".jpg":      _from_image,
    ".jpeg":     _from_image,
    ".png":      _from_image,
    ".bmp":      _from_image,
    ".tiff":     _from_image,
    ".tif":      _from_image,
    ".gif":      _from_image,
    ".webp":     _from_image,
}

# Public list for UI / validation
SUPPORTED_EXTENSIONS = set(_EXTRACTORS.keys())

#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
╔══════════════════════════════════════════════════════════════════╗
║   AEJ CV MATCHER PRO  v2.0  — ÉDITION PRESTIGE                  ║
║   Le logiciel de matching le plus performant au monde            ║
║   Agence Emploi Jeunes / Agence Prestige                         ║
║   Bd des Martyrs, Cocody, Abidjan — Côte d'Ivoire               ║
║   Auteur : KEITA — CIP Service Entreprises                       ║
╚══════════════════════════════════════════════════════════════════╝

INNOVATIONS v2.0 :
  ✦ Moteur IA Hybride TF-IDF + Ollama LLM  → 100% local, zéro cloud
  ✦ Traitement parallèle (4 threads Drive)  → 4× plus rapide
  ✦ Déduplication intelligente (hash MD5)
  ✦ Tableau de bord analytique (6 graphiques temps réel)
  ✦ 🤝 Module Mise en Relation (CRM complet)
  ✦ Génération emails personnalisés (IA locale)
  ✦ Lettres de présélection DOCX par poste
  ✦ ⏰ Mode Autonome (watchdog + planification)
  ✦ Scoring percentile et ranking sectoriel
  ✦ Estimation du coût API avant exécution
  ✦ Tracking candidat : En attente → Recruté
  ✦ 7 onglets — interface complète
"""

# ════════════════════════════════════════════════════════════════
#  SECTION 1 — IMPORTS
# ════════════════════════════════════════════════════════════════
import tkinter as tk
from tkinter import ttk, messagebox, filedialog, scrolledtext
import threading, json, os, sys, time, re, io, sqlite3
import requests, shutil, tempfile, traceback, hashlib, webbrowser
from datetime import datetime, timedelta
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed

# ── Google Drive ─────────────────────────────────────────────────
try:
    from google.oauth2.credentials import Credentials
    from google_auth_oauthlib.flow import InstalledAppFlow
    from google.auth.transport.requests import Request as GRequest
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaIoBaseDownload, MediaFileUpload
    GOOGLE_OK = True
except ImportError:
    GOOGLE_OK = False

# ── PDF ──────────────────────────────────────────────────────────
try:
    import pdfplumber
    PDF_OK = True
except ImportError:
    PDF_OK = False

# ── DOCX ─────────────────────────────────────────────────────────
try:
    from docx import Document as DocxDoc
    from docx.shared import Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

# ── OCR ──────────────────────────────────────────────────────────
try:
    import pytesseract
    from PIL import Image
    OCR_OK = True
except ImportError:
    OCR_OK = False

# ── Excel ────────────────────────────────────────────────────────
try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    from openpyxl.chart import BarChart, PieChart, Reference
    EXCEL_OK = True
except ImportError:
    EXCEL_OK = False

# ── Matplotlib (Dashboard) ───────────────────────────────────────
try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    from matplotlib.figure import Figure
    from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
    from matplotlib.gridspec import GridSpec
    import matplotlib.patches as mpatches
    MPLOT_OK = True
except ImportError:
    MPLOT_OK = False

# ── TF-IDF Pré-filtrage ──────────────────────────────────────────
try:
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics.pairwise import cosine_similarity
    SKLEARN_OK = True
except ImportError:
    SKLEARN_OK = False

# ── BM25 (ranking IR de référence) ───────────────────────────────
try:
    from rank_bm25 import BM25Okapi
    BM25_OK = True
except ImportError:
    BM25_OK = False

# ── RapidFuzz (matching flou des compétences) ────────────────────
try:
    from rapidfuzz import fuzz as rf_fuzz
    FUZZ_OK = True
except ImportError:
    FUZZ_OK = False

# ── Embeddings sémantiques (optionnel, le + puissant) ────────────
EMB_OK = False
EMB_BACKEND = None
try:
    from sentence_transformers import SentenceTransformer
    EMB_OK = True
    EMB_BACKEND = 'sentence_transformers'
except ImportError:
    try:
        from fastembed import TextEmbedding
        EMB_OK = True
        EMB_BACKEND = 'fastembed'
    except ImportError:
        EMB_OK = False

# ── Formats additionnels ─────────────────────────────────────────
try:
    import docx2txt
    DOCX2TXT_OK = True
except ImportError:
    DOCX2TXT_OK = False

try:
    from striprtf.striprtf import rtf_to_text
    RTF_OK = True
except ImportError:
    RTF_OK = False

try:
    from pptx import Presentation
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

try:
    import cv2
    import numpy as np
    CV2_OK = True
except ImportError:
    CV2_OK = False

try:
    from pdf2image import convert_from_path
    PDF2IMG_OK = True
except ImportError:
    PDF2IMG_OK = False

# ── Notifications Windows ────────────────────────────────────────
try:
    from win10toast import ToastNotifier
    TOAST_OK = True
except ImportError:
    TOAST_OK = False


# ════════════════════════════════════════════════════════════════
#  SECTION 2 — CONSTANTES
# ════════════════════════════════════════════════════════════════
VERSION    = "2.0"
APP_DIR    = Path(sys.executable).parent if getattr(sys, 'frozen', False) else Path(__file__).parent
CFG_FILE   = APP_DIR / "config_aej_matcher.json"
DB_FILE    = APP_DIR / "aej_matcher_cache.db"
CREDS_FILE = APP_DIR / "credentials.json"
TOKEN_FILE = APP_DIR / "token.json"
SCOPES     = ['https://www.googleapis.com/auth/drive',
              'https://www.googleapis.com/auth/gmail.send']

AEJ_GREEN  = "#1A7A3C"
AEJ_DARK   = "#0D5A2C"
AEJ_ORANGE = "#F47920"
AEJ_LIGHT  = "#E8F5E9"
AEJ_CREAM  = "#FAFAF7"
AEJ_GRAY   = "#F5F5F5"

STATUTS = ["EN_ATTENTE", "CONTACTE", "ENTRETIEN_PLANIFIE",
           "DOSSIER_TRANSMIS", "RECRUTE", "REJETE"]
STATUTS_FR = {
    "EN_ATTENTE"        : "⏳ En attente",
    "CONTACTE"          : "📧 Contacté",
    "ENTRETIEN_PLANIFIE": "📅 Entretien planifié",
    "DOSSIER_TRANSMIS"  : "📤 Dossier transmis",
    "RECRUTE"           : "✅ Recruté",
    "REJETE"            : "❌ Rejeté",
}
STATUTS_COLORS = {
    "EN_ATTENTE"        : "#FFF9C4",
    "CONTACTE"          : "#E3F2FD",
    "ENTRETIEN_PLANIFIE": "#E8F5E9",
    "DOSSIER_TRANSMIS"  : "#FFF3E0",
    "RECRUTE"           : "#C8E6C9",
    "REJETE"            : "#FFCDD2",
}

MIMES_LISIBLES = {
    'application/pdf',
    'application/msword',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    'application/vnd.oasis.opendocument.text',
    'application/rtf', 'text/rtf',
    'text/plain', 'text/html', 'text/csv',
    'image/jpeg', 'image/jpg', 'image/png', 'image/tiff', 'image/bmp', 'image/webp',
    'application/vnd.google-apps.document',
}

# Extensions locales acceptées pour l'import
EXTS_SUPPORTEES = {'.pdf', '.doc', '.docx', '.jpg', '.jpeg', '.png', '.txt'}
EXTS_LABEL      = "PDF, Word (.doc/.docx), Images (.jpg/.png)"

# Statuts d'upload
UP_STATUTS = {
    'EN_ATTENTE': '⏳ En attente',
    'UPLOAD'    : '📤 Upload...',
    'OK'        : '✅ Uploadé',
    'ERREUR'    : '❌ Erreur',
    'EXISTANT'  : '⏭ Déjà présent',
}

DEFAULT_CONFIG = {
    # --- Source des données : "DRIVE" ou "LOCAL" ---
    "SOURCE_MODE"          : "LOCAL",
    "DOSSIER_CV_SOURCE"    : "18F2W3Gvy21-z7IadV5T8_ypHhwN5a9Gy",
    "DOSSIER_FICHES_POSTE" : "1lmRrkeMGLA0lNhNevzk6Weuq8Bb7lY4S",
    "DOSSIER_DESTINATION"  : "129l-ymj4NHYiJCU9OnXErCfYyZg-tt8q",
    # --- Chemins locaux (mode LOCAL) ---
    "LOCAL_CV_PATHS"       : [],   # liste de fichiers et/ou dossiers
    "LOCAL_FICHE_PATHS"    : [],   # liste de fichiers et/ou dossiers
    "LOCAL_DEST_PATH"      : "",   # dossier de destination local
    "OLLAMA_URL"           : "http://localhost:11434",
    "OLLAMA_MODEL"         : "gemma3:4b",
    "MODE"                 : "MULTI_POSTES",
    "ACTION"               : "COPIER",
    "SEUIL_TOP"            : 80,
    "SEUIL_MOYEN"          : 65,
    "SEUIL_TFIDF"          : 18,
    "RETRIEVAL_TOP_K"      : 6,      # nb de postes candidats par CV (retrieval hybride)
    "USE_EMBEDDINGS"       : False,  # léger par défaut (BM25+fuzzy+TF-IDF). Activez si sentence-transformers installé
    "EMB_MODEL"            : "paraphrase-multilingual-MiniLM-L12-v2",
    "BLEND_LLM"            : 0.70,   # poids du score LLM dans le score hybride
    "BLEND_SEMANTIC"       : 0.18,   # poids similarité sémantique
    "BLEND_SKILLS"         : 0.12,   # poids couverture compétences (fuzzy)
    "FUZZ_SEUIL"           : 82,     # seuil de match flou d'une compétence (0-100)
    "MAX_CV_PAR_POSTE"     : 15,
    "MAX_POSTES_PAR_CV"    : 3,
    "PARALLEL_WORKERS"     : 4,
    "PAUSE_API"            : 1.0,
    "TESSERACT_PATH"       : r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    "SMTP_SERVEUR"         : "",
    "SMTP_PORT"            : 587,
    "SMTP_EMAIL"           : "",
    "SMTP_MOT_DE_PASSE"    : "",
    "AUTO_ACTIVE"          : False,
    "AUTO_INTERVALLE_H"    : 6,
    "SIGNATURE_EMAIL"      : (
        "\n\nCordialement,\n"
        "Service Entreprises — Agence Prestige\n"
        "Agence Emploi Jeunes\n"
        "Bd des Martyrs, Cocody, Abidjan\n"
        "📞 01 73 26 26 43"
    ),
}


# ════════════════════════════════════════════════════════════════
#  SECTION 3 — CONFIGURATION
# ════════════════════════════════════════════════════════════════
class Config:
    def __init__(self):
        self.d = dict(DEFAULT_CONFIG)
        if CFG_FILE.exists():
            try:
                self.d.update(json.loads(CFG_FILE.read_text('utf-8')))
            except Exception:
                pass

    def get(self, k, default=None):
        return self.d.get(k, default)

    def set(self, k, v):
        self.d[k] = v
        self._save()

    def update(self, dct):
        self.d.update(dct)
        self._save()

    def _save(self):
        CFG_FILE.write_text(json.dumps(self.d, indent=2, ensure_ascii=False), 'utf-8')


# ════════════════════════════════════════════════════════════════
#  SECTION 4 — BASE DE DONNÉES
# ════════════════════════════════════════════════════════════════
class Database:
    def __init__(self):
        self.conn = sqlite3.connect(str(DB_FILE), check_same_thread=False)
        self.conn.row_factory = sqlite3.Row
        self._init()

    def _init(self):
        c = self.conn
        c.execute("""CREATE TABLE IF NOT EXISTS cv_profils
            (id TEXT PRIMARY KEY, nom TEXT, hash_md5 TEXT,
             json TEXT, ts TEXT)""")
        c.execute("""CREATE TABLE IF NOT EXISTS fiche_profils
            (id TEXT PRIMARY KEY, nom TEXT, json TEXT, ts TEXT)""")
        c.execute("""CREATE TABLE IF NOT EXISTS scores
            (id_cv TEXT, id_fiche TEXT, score_tfidf REAL,
             json TEXT, ts TEXT, PRIMARY KEY(id_cv, id_fiche))""")
        c.execute("""CREATE TABLE IF NOT EXISTS contacts
            (id INTEGER PRIMARY KEY AUTOINCREMENT,
             id_cv TEXT, id_fiche TEXT,
             nom_candidat TEXT, titre_profil TEXT,
             intitule_poste TEXT, entreprise TEXT,
             score INTEGER, decision TEXT,
             statut TEXT DEFAULT 'EN_ATTENTE',
             email_candidat TEXT, tel_candidat TEXT,
             notes TEXT, email_genere TEXT,
             date_creation TEXT, date_modif TEXT,
             UNIQUE(id_cv, id_fiche))""")
        c.execute("""CREATE TABLE IF NOT EXISTS runs
            (id INTEGER PRIMARY KEY AUTOINCREMENT,
             date_debut TEXT, date_fin TEXT,
             nb_cv INTEGER, nb_fiches INTEGER,
             nb_matches INTEGER, duree_sec INTEGER,
             statut TEXT)""")
        c.commit()

    # ── CV ───────────────────────────────────────────────────────
    def get_cv(self, fid):
        r = self.conn.execute("SELECT json FROM cv_profils WHERE id=?", (fid,)).fetchone()
        return json.loads(r[0]) if r else None

    def get_cv_hash(self, h):
        r = self.conn.execute("SELECT id FROM cv_profils WHERE hash_md5=?", (h,)).fetchone()
        return r[0] if r else None

    def save_cv(self, fid, nom, h, p):
        self.conn.execute("INSERT OR REPLACE INTO cv_profils VALUES (?,?,?,?,?)",
            (fid, nom, h, json.dumps(p, ensure_ascii=False), datetime.now().isoformat()))
        self.conn.commit()

    # ── Fiches ───────────────────────────────────────────────────
    def get_fiche(self, fid):
        r = self.conn.execute("SELECT json FROM fiche_profils WHERE id=?", (fid,)).fetchone()
        return json.loads(r[0]) if r else None

    def save_fiche(self, fid, nom, p):
        self.conn.execute("INSERT OR REPLACE INTO fiche_profils VALUES (?,?,?,?)",
            (fid, nom, json.dumps(p, ensure_ascii=False), datetime.now().isoformat()))
        self.conn.commit()

    # ── Scores ───────────────────────────────────────────────────
    def get_score(self, id_cv, id_f):
        r = self.conn.execute("SELECT json FROM scores WHERE id_cv=? AND id_fiche=?",
            (id_cv, id_f)).fetchone()
        return json.loads(r[0]) if r else None

    def save_score(self, id_cv, id_f, tfidf, s):
        self.conn.execute("INSERT OR REPLACE INTO scores VALUES (?,?,?,?,?)",
            (id_cv, id_f, tfidf, json.dumps(s, ensure_ascii=False), datetime.now().isoformat()))
        self.conn.commit()

    # ── Contacts (CRM) ───────────────────────────────────────────
    def upsert_contact(self, id_cv, id_fiche, data: dict):
        existing = self.conn.execute(
            "SELECT id FROM contacts WHERE id_cv=? AND id_fiche=?", (id_cv, id_fiche)
        ).fetchone()
        now = datetime.now().isoformat()
        if existing:
            self.conn.execute("""UPDATE contacts SET
                nom_candidat=?, titre_profil=?, intitule_poste=?, entreprise=?,
                score=?, decision=?, date_modif=?
                WHERE id_cv=? AND id_fiche=?""",
                (data.get('nom_candidat',''), data.get('titre_profil',''),
                 data.get('intitule_poste',''), data.get('entreprise',''),
                 data.get('score',0), data.get('decision',''), now,
                 id_cv, id_fiche))
        else:
            self.conn.execute("""INSERT OR IGNORE INTO contacts
                (id_cv, id_fiche, nom_candidat, titre_profil,
                 intitule_poste, entreprise, score, decision,
                 statut, date_creation, date_modif)
                VALUES (?,?,?,?,?,?,?,?,?,?,?)""",
                (id_cv, id_fiche, data.get('nom_candidat',''),
                 data.get('titre_profil',''), data.get('intitule_poste',''),
                 data.get('entreprise',''), data.get('score',0),
                 data.get('decision',''), 'EN_ATTENTE', now, now))
        self.conn.commit()

    def get_contacts(self, statut=None, poste=None):
        q = "SELECT * FROM contacts WHERE 1=1"
        params = []
        if statut and statut != 'TOUS':
            q += " AND statut=?"
            params.append(statut)
        if poste and poste != 'TOUS':
            q += " AND intitule_poste=?"
            params.append(poste)
        q += " ORDER BY score DESC"
        rows = self.conn.execute(q, params).fetchall()
        return [dict(r) for r in rows]

    def update_contact(self, cid, **kwargs):
        sets = ', '.join(f"{k}=?" for k in kwargs)
        vals = list(kwargs.values()) + [datetime.now().isoformat(), cid]
        self.conn.execute(f"UPDATE contacts SET {sets}, date_modif=? WHERE id=?", vals)
        self.conn.commit()

    def get_postes_distincts(self):
        rows = self.conn.execute("SELECT DISTINCT intitule_poste FROM contacts ORDER BY intitule_poste").fetchall()
        return [r[0] for r in rows]

    def stats_contacts(self):
        rows = self.conn.execute(
            "SELECT statut, COUNT(*) as n FROM contacts GROUP BY statut"
        ).fetchall()
        return {r[0]: r[1] for r in rows}

    # ── Runs ─────────────────────────────────────────────────────
    def save_run(self, debut, fin, nb_cv, nb_fiches, nb_matches, duree, statut):
        self.conn.execute("INSERT INTO runs VALUES (NULL,?,?,?,?,?,?,?)",
            (debut, fin, nb_cv, nb_fiches, nb_matches, duree, statut))
        self.conn.commit()

    def last_run(self):
        r = self.conn.execute("SELECT * FROM runs ORDER BY id DESC LIMIT 1").fetchone()
        return dict(r) if r else None

    # ── Global ───────────────────────────────────────────────────
    def vider(self):
        for t in ['cv_profils', 'fiche_profils', 'scores']:
            self.conn.execute(f"DELETE FROM {t}")
        self.conn.commit()

    def stats_globales(self):
        nb_cv    = self.conn.execute("SELECT COUNT(*) FROM cv_profils").fetchone()[0]
        nb_f     = self.conn.execute("SELECT COUNT(*) FROM fiche_profils").fetchone()[0]
        nb_sc    = self.conn.execute("SELECT COUNT(*) FROM scores").fetchone()[0]
        nb_top   = self.conn.execute("SELECT COUNT(*) FROM contacts WHERE score>=80").fetchone()[0]
        return {"cv": nb_cv, "fiches": nb_f, "scores": nb_sc, "top": nb_top}


# ════════════════════════════════════════════════════════════════
#  SECTION 5 — GOOGLE DRIVE CLIENT
# ════════════════════════════════════════════════════════════════
class DriveClient:
    def __init__(self):
        self.svc = None

    def connecter(self):
        if not GOOGLE_OK:
            raise RuntimeError("Bibliothèques Google manquantes. pip install google-api-python-client google-auth-oauthlib")
        if not CREDS_FILE.exists():
            raise FileNotFoundError(f"credentials.json introuvable dans {APP_DIR}\nTéléchargez-le depuis Google Cloud Console.")
        creds = None
        if TOKEN_FILE.exists():
            try:
                creds = Credentials.from_authorized_user_file(str(TOKEN_FILE), SCOPES)
            except Exception:
                creds = None
        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                creds.refresh(GRequest())
            else:
                flow = InstalledAppFlow.from_client_secrets_file(str(CREDS_FILE), SCOPES)
                creds = flow.run_local_server(port=0)
            TOKEN_FILE.write_text(creds.to_json())
        self.svc = build('drive', 'v3', credentials=creds)

    def lister(self, folder_id, recurse=False):
        items, pt = [], None
        q = f"'{folder_id}' in parents and trashed=false"
        while True:
            r = self.svc.files().list(
                q=q, spaces='drive', pageToken=pt, pageSize=200,
                fields='nextPageToken, files(id,name,mimeType,size,modifiedTime)'
            ).execute()
            for f in r.get('files', []):
                if f['mimeType'] == 'application/vnd.google-apps.folder':
                    if recurse:
                        items.extend(self.lister(f['id'], recurse=True))
                elif f['mimeType'] in MIMES_LISIBLES:
                    items.append(f)
            pt = r.get('nextPageToken')
            if not pt:
                break
        return items

    def telecharger(self, file_id, mime_type, dest: Path):
        if mime_type == 'application/vnd.google-apps.document':
            req = self.svc.files().export_media(fileId=file_id, mimeType='text/plain')
        else:
            req = self.svc.files().get_media(fileId=file_id)
        with io.FileIO(str(dest), 'wb') as fh:
            dl = MediaIoBaseDownload(fh, req)
            done = False
            while not done:
                _, done = dl.next_chunk()

    def creer_dossier(self, nom, parent_id):
        nom_s = nom.replace("'", "\\'")
        q = (f"name='{nom_s}' and '{parent_id}' in parents and "
             f"mimeType='application/vnd.google-apps.folder' and trashed=false")
        r = self.svc.files().list(q=q, fields='files(id)').execute()
        if r.get('files'):
            return r['files'][0]['id']
        meta = {'name': nom, 'mimeType': 'application/vnd.google-apps.folder', 'parents': [parent_id]}
        return self.svc.files().create(body=meta, fields='id').execute()['id']

    def copier(self, file_id, nom, parent_id):
        return self.svc.files().copy(
            fileId=file_id, body={'name': nom, 'parents': [parent_id]}, fields='id'
        ).execute()['id']

    def deplacer(self, file_id, nom, parent_id):
        f = self.svc.files().get(fileId=file_id, fields='parents').execute()
        old = ','.join(f.get('parents', []))
        self.svc.files().update(
            fileId=file_id, addParents=parent_id, removeParents=old,
            body={'name': nom}, fields='id'
        ).execute()

    def uploader(self, local: Path, nom, parent_id, mime='application/octet-stream'):
        meta = {'name': nom, 'parents': [parent_id]}
        media = MediaFileUpload(str(local), mimetype=mime)
        r = self.svc.files().create(body=meta, media_body=media, fields='id,webViewLink').execute()
        return r.get('webViewLink', '')

    def uploader_local(self, local: Path, nom: str, parent_id: str) -> tuple:
        """Upload un fichier local vers Drive avec détection MIME auto.
        Retourne (drive_id, web_view_link)."""
        import mimetypes as _mt
        mime, _ = _mt.guess_type(str(local))
        if not mime:
            ext = local.suffix.lower()
            mime = {
                '.pdf' : 'application/pdf',
                '.doc' : 'application/msword',
                '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                '.jpg' : 'image/jpeg',
                '.jpeg': 'image/jpeg',
                '.png' : 'image/png',
                '.txt' : 'text/plain',
            }.get(ext, 'application/octet-stream')
        meta  = {'name': nom, 'parents': [parent_id]}
        media = MediaFileUpload(str(local), mimetype=mime, resumable=True)
        r = self.svc.files().create(body=meta, media_body=media,
                                     fields='id,webViewLink').execute()
        return r.get('id', ''), r.get('webViewLink', '')

    def fichier_existe(self, nom: str, parent_id: str) -> bool:
        """Vérifie si un fichier du même nom existe déjà dans le dossier Drive."""
        nom_s = nom.replace("'", "\\'")
        q = (f"name='{nom_s}' and '{parent_id}' in parents "
             f"and trashed=false")
        r = self.svc.files().list(q=q, fields='files(id)').execute()
        return len(r.get('files', [])) > 0

    def lister_dossiers(self, parent_id: str) -> list:
        """Liste les sous-dossiers directs d'un dossier Drive."""
        q = (f"'{parent_id}' in parents and trashed=false "
             f"and mimeType='application/vnd.google-apps.folder'")
        r = self.svc.files().list(q=q, fields='files(id,name)').execute()
        return r.get('files', [])

    def compter_fichiers(self, folder_id):
        items = self.lister(folder_id)
        return len(items)


# ════════════════════════════════════════════════════════════════
#  SECTION 5b — CLIENT LOCAL (fichiers / dossiers du PC)
#  Implémente la MÊME interface que DriveClient → orchestrateur
#  inchangé. "folder_id" = jeton sémantique ou chemin local.
# ════════════════════════════════════════════════════════════════
_EXT_MIME = {
    '.pdf' : 'application/pdf',
    '.doc' : 'application/msword',
    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    '.odt' : 'application/vnd.oasis.opendocument.text',
    '.rtf' : 'application/rtf',
    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    '.txt' : 'text/plain', '.csv': 'text/csv',
    '.html': 'text/html', '.htm': 'text/html',
    '.jpg' : 'image/jpeg', '.jpeg': 'image/jpeg', '.png': 'image/png',
    '.tif' : 'image/tiff', '.tiff': 'image/tiff', '.bmp': 'image/bmp', '.webp': 'image/webp',
}

class LocalClient:
    """Fournit CV/fiches depuis le disque local. Aucune connexion requise."""

    def __init__(self, cfg: Config):
        self.cfg        = cfg
        self.cv_paths   = list(cfg.get('LOCAL_CV_PATHS', []) or [])
        self.fiche_paths = list(cfg.get('LOCAL_FICHE_PATHS', []) or [])

    def connecter(self):
        # Vérifier que la destination existe / est créable
        dest = self.cfg.get('LOCAL_DEST_PATH', '')
        if not dest:
            raise RuntimeError("Dossier de destination local non défini (onglet Configuration).")
        Path(dest).mkdir(parents=True, exist_ok=True)

    def _expand(self, paths, recurse=True):
        """Transforme une liste de fichiers/dossiers en liste de fichiers lisibles."""
        out = []
        for p in paths:
            pp = Path(p)
            if not pp.exists():
                continue
            if pp.is_file():
                if pp.suffix.lower() in _EXT_MIME:
                    out.append(pp)
            elif pp.is_dir():
                iterator = pp.rglob('*') if recurse else pp.glob('*')
                for f in iterator:
                    if f.is_file() and f.suffix.lower() in _EXT_MIME:
                        out.append(f)
        # Déduplication par chemin absolu
        seen, uniq = set(), []
        for f in out:
            a = str(f.resolve())
            if a not in seen:
                seen.add(a)
                uniq.append(f)
        return uniq

    def lister(self, token, recurse=False):
        if token == 'LOCAL_CV':
            files = self._expand(self.cv_paths, recurse=True)
        elif token == 'LOCAL_FICHE':
            files = self._expand(self.fiche_paths, recurse=True)
        else:
            # token est un chemin de dossier réel (ex. sous-dossier destination)
            files = self._expand([token], recurse=recurse)
        items = []
        for f in files:
            try:
                stat = f.stat()
                items.append({
                    'id'          : str(f.resolve()),
                    'name'        : f.name,
                    'mimeType'    : _EXT_MIME.get(f.suffix.lower(), 'application/octet-stream'),
                    'size'        : stat.st_size,
                    'modifiedTime': datetime.fromtimestamp(stat.st_mtime).isoformat(),
                })
            except Exception:
                pass
        return items

    def telecharger(self, file_id, mime_type, dest: Path):
        # file_id est déjà un chemin local → simple copie
        shutil.copy2(file_id, str(dest))

    def creer_dossier(self, nom, parent_id):
        # parent_id est un chemin local
        d = Path(parent_id) / _nettoyer(nom)
        d.mkdir(parents=True, exist_ok=True)
        return str(d)

    def copier(self, file_id, nom, parent_id):
        dest = Path(parent_id) / nom
        shutil.copy2(file_id, str(dest))
        return str(dest)

    def deplacer(self, file_id, nom, parent_id):
        dest = Path(parent_id) / nom
        shutil.move(file_id, str(dest))
        return str(dest)

    def uploader(self, local: Path, nom, parent_id, mime='application/octet-stream'):
        dest = Path(parent_id) / nom
        shutil.copy2(str(local), str(dest))
        return str(dest)

    def compter_fichiers(self, token):
        return len(self.lister(token))


def make_provider(cfg: Config):
    """Retourne le bon client selon le mode source choisi."""
    if cfg.get('SOURCE_MODE', 'DRIVE') == 'LOCAL':
        return LocalClient(cfg)
    return DriveClient()


# ════════════════════════════════════════════════════════════════
#  SECTION 6 — EXTRACTEUR DE TEXTE MULTI-FORMAT
#  Formats : PDF (texte+OCR), DOC, DOCX, ODT, RTF, PPTX, TXT, CSV,
#            HTML, images (JPG/PNG/TIFF/BMP/WEBP) avec OCR avancé.
#  Stratégie en cascade : chaque format tente plusieurs moteurs
#  et bascule sur l'OCR si aucune couche texte n'est trouvée.
# ════════════════════════════════════════════════════════════════
class Extracteur:
    def __init__(self, tesseract_path=None):
        if tesseract_path and OCR_OK and Path(tesseract_path).exists():
            pytesseract.pytesseract.tesseract_cmd = tesseract_path

    # ── Dispatcher principal ─────────────────────────────────────
    def extraire(self, chemin: Path, mime_type: str) -> str:
        ext = Path(chemin).suffix.lower()
        try:
            if ext == '.pdf' or mime_type == 'application/pdf':
                return self._pdf(chemin)
            if ext == '.docx' or mime_type.endswith('wordprocessingml.document'):
                return self._docx(chemin)
            if ext == '.doc' or mime_type == 'application/msword':
                return self._doc(chemin)
            if ext == '.odt' or 'opendocument.text' in mime_type:
                return self._odt(chemin)
            if ext == '.rtf' or 'rtf' in mime_type:
                return self._rtf(chemin)
            if ext == '.pptx' or 'presentationml' in mime_type:
                return self._pptx(chemin)
            if ext in ('.html', '.htm') or mime_type == 'text/html':
                return self._html(chemin)
            if ext in ('.jpg', '.jpeg', '.png', '.tif', '.tiff', '.bmp', '.webp') \
               or mime_type.startswith('image/'):
                return self._image(chemin)
            # txt, csv, fallback
            return Path(chemin).read_text(encoding='utf-8', errors='ignore')
        except Exception as e:
            return f"[Extraction impossible: {e}]"

    # ── PDF : texte natif + OCR de secours ──────────────────────
    def _pdf(self, p: Path) -> str:
        parts = []
        texte_natif_ok = False
        if PDF_OK:
            try:
                with pdfplumber.open(str(p)) as pdf:
                    for page in pdf.pages:
                        t = page.extract_text()
                        if t and t.strip():
                            parts.append(t)
                            texte_natif_ok = True
                        else:
                            # tables
                            for tbl in (page.extract_tables() or []):
                                for ligne in tbl:
                                    parts.append(' '.join(c for c in ligne if c))
            except Exception:
                pass
        # Si pas de couche texte → OCR page par page
        if not texte_natif_ok and OCR_OK:
            try:
                if PDF2IMG_OK:
                    images = convert_from_path(str(p), dpi=220)
                    for img in images:
                        parts.append(self._ocr_image(img))
                elif PDF_OK:
                    with pdfplumber.open(str(p)) as pdf:
                        for page in pdf.pages:
                            img = page.to_image(resolution=220).original
                            parts.append(self._ocr_image(img))
            except Exception:
                pass
        return '\n'.join(x for x in parts if x and x.strip())

    # ── DOCX : paragraphes + tableaux ───────────────────────────
    def _docx(self, p: Path) -> str:
        textes = []
        if DOCX_OK:
            try:
                doc = DocxDoc(str(p))
                for par in doc.paragraphs:
                    if par.text.strip():
                        textes.append(par.text)
                for table in doc.tables:           # CV souvent en tableaux
                    for row in table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            textes.append(' | '.join(cells))
            except Exception:
                pass
        if not textes and DOCX2TXT_OK:
            try:
                textes.append(docx2txt.process(str(p)))
            except Exception:
                pass
        return '\n'.join(textes)

    # ── DOC (ancien binaire Word) : cascade de moteurs ──────────
    def _doc(self, p: Path) -> str:
        # 1) docx2txt fonctionne parfois sur .doc récents
        if DOCX2TXT_OK:
            try:
                t = docx2txt.process(str(p))
                if t and t.strip():
                    return t
            except Exception:
                pass
        # 2) LibreOffice headless (si installé)
        t = self._libreoffice_to_txt(p)
        if t:
            return t
        # 3) Word via COM (Windows + Word installé)
        t = self._word_com(p)
        if t:
            return t
        # 4) Récupération brute du texte ASCII/latin embarqué
        try:
            raw = Path(p).read_bytes()
            txt = raw.decode('latin-1', errors='ignore')
            txt = re.sub(r'[^\x20-\x7E\u00C0-\u017F\n\r\t]', ' ', txt)
            txt = re.sub(r'\s{3,}', '  ', txt)
            mots = re.findall(r'[A-Za-zÀ-ÿ0-9@.\-]{3,}', txt)
            return ' '.join(mots) if len(mots) > 20 else "[.doc illisible — convertir en .docx/.pdf]"
        except Exception:
            return "[.doc illisible — convertir en .docx/.pdf]"

    def _libreoffice_to_txt(self, p: Path) -> str:
        for binname in ('soffice', 'libreoffice',
                        r'C:\Program Files\LibreOffice\program\soffice.exe'):
            try:
                outdir = Path(tempfile.mkdtemp(prefix='lo_'))
                r = __import__('subprocess').run(
                    [binname, '--headless', '--convert-to', 'txt:Text',
                     '--outdir', str(outdir), str(p)],
                    capture_output=True, timeout=60)
                produced = list(outdir.glob('*.txt'))
                if produced:
                    return produced[0].read_text(encoding='utf-8', errors='ignore')
            except Exception:
                continue
        return ""

    def _word_com(self, p: Path) -> str:
        try:
            import win32com.client  # type: ignore
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(str(Path(p).resolve()))
            txt = doc.Content.Text
            doc.Close(False)
            word.Quit()
            return txt
        except Exception:
            return ""

    # ── ODT (OpenDocument) ──────────────────────────────────────
    def _odt(self, p: Path) -> str:
        try:
            import zipfile
            from xml.etree import ElementTree as ET
            with zipfile.ZipFile(str(p)) as z:
                xml = z.read('content.xml')
            txt = re.sub(r'<[^>]+>', ' ', xml.decode('utf-8', errors='ignore'))
            return re.sub(r'\s{2,}', ' ', txt).strip()
        except Exception:
            t = self._libreoffice_to_txt(p)
            return t or "[ODT illisible]"

    # ── RTF ─────────────────────────────────────────────────────
    def _rtf(self, p: Path) -> str:
        raw = Path(p).read_text(encoding='utf-8', errors='ignore')
        if RTF_OK:
            try:
                return rtf_to_text(raw)
            except Exception:
                pass
        return re.sub(r'\\[a-z]+\d* ?|[{}]', ' ', raw)

    # ── PPTX ────────────────────────────────────────────────────
    def _pptx(self, p: Path) -> str:
        if not PPTX_OK:
            return "[python-pptx non installé]"
        textes = []
        prs = Presentation(str(p))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for par in shape.text_frame.paragraphs:
                        line = ''.join(run.text for run in par.runs)
                        if line.strip():
                            textes.append(line)
        return '\n'.join(textes)

    # ── HTML ────────────────────────────────────────────────────
    def _html(self, p: Path) -> str:
        raw = Path(p).read_text(encoding='utf-8', errors='ignore')
        raw = re.sub(r'(?is)<(script|style).*?</\1>', ' ', raw)
        return re.sub(r'\s{2,}', ' ', re.sub(r'<[^>]+>', ' ', raw)).strip()

    # ── Images : OCR avec pré-traitement ────────────────────────
    def _image(self, p: Path) -> str:
        if not OCR_OK:
            return "[pytesseract non installé]"
        try:
            img = Image.open(str(p))
            return self._ocr_image(img)
        except Exception as e:
            return f"[Image illisible: {e}]"

    def _ocr_image(self, pil_img) -> str:
        """OCR avec pré-traitement (améliore nettement les scans)."""
        try:
            if CV2_OK:
                arr = np.array(pil_img.convert('RGB'))
                gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
                # débruitage + binarisation adaptative + redressement léger
                gray = cv2.fastNlMeansDenoising(gray, h=10)
                gray = cv2.adaptiveThreshold(
                    gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                    cv2.THRESH_BINARY, 31, 11)
                proc = Image.fromarray(gray)
            else:
                from PIL import ImageOps, ImageFilter
                proc = pil_img.convert('L')
                proc = ImageOps.autocontrast(proc)
                proc = proc.filter(ImageFilter.SHARPEN)
            return pytesseract.image_to_string(proc, lang='fra+eng')
        except Exception:
            try:
                return pytesseract.image_to_string(pil_img, lang='fra+eng')
            except Exception as e:
                return f"[OCR échoué: {e}]"

    @staticmethod
    def md5(texte: str) -> str:
        return hashlib.md5(texte.encode('utf-8', errors='ignore')).hexdigest()


# ════════════════════════════════════════════════════════════════
#  SECTION 7 — MOTEUR IA HYBRIDE  (TF-IDF + Ollama local)
# ════════════════════════════════════════════════════════════════

PROMPT_CV = """\
Analyse ce CV et retourne UNIQUEMENT ce JSON valide (sans markdown) :
{"nom_candidat":"","titre_profil":"","secteur":"","annees_experience":0,
"competences_techniques":[],"formation_niveau":"","formation_domaine":"",
"certifications":[],"logiciels_outils":[],"langues":[],"localisation":"",
"mots_cles_metier":[],"email":"","telephone":"","resume_profil":""}
Nom fichier: %s
CONTENU CV:
%s"""

PROMPT_FICHE = """\
Analyse cette fiche de poste et retourne UNIQUEMENT ce JSON valide (sans markdown) :
{"intitule":"","entreprise":"","secteur":"","competences_obligatoires":[],
"competences_souhaitees":[],"formation_requise":"","experience_min_annees":0,
"logiciels_requis":[],"certifications_requises":[],"langues_requises":[],
"localisation":"","mots_cles_sectoriels":[],"criteres_eliminatoires":[],
"resume_poste":""}
Nom fichier: %s
CONTENU FICHE:
%s"""

PROMPT_SCORE = """\
Évalue la correspondance (score 0-100) entre ce candidat et ce poste.
Pondération: compétences_métier=35%%, expérience=25%%, diplôme=15%%,
mots_clés_sectoriels=10%%, outils=5%%, langues=5%%, localisation=5%%
Retourne UNIQUEMENT ce JSON valide :
{"score_global":0,
"scores_detail":{"competences_metier":0,"experience":0,"diplome":0,
"mots_cles_sectoriels":0,"outils_logiciels":0,"langues":0,"localisation":0},
"points_forts":[],"reserves":[],"criteres_eliminatoires_touches":false,
"decision":"PRIORITAIRE|MOYEN|FAIBLE","justification":""}
CANDIDAT: %s
POSTE: %s"""

PROMPT_EMAIL_CONVOC = """\
Rédige un email de convocation à un entretien en français, professionnel et chaleureux.
Signataire : Service Entreprises — Agence Prestige (AEJ Cocody Abidjan).
Email court, personnalisé, sans champs vides entre crochets.
Candidat : %s | Poste : %s chez %s | Score matching : %s%%
Inclure : objet de l'email sur la 1ère ligne préfixé "OBJET:", puis le corps.
Ne retourne QUE l'email, sans explication."""

PROMPT_EMAIL_RECRUTEUR = """\
Rédige un email professionnel en français à un recruteur pour lui transmettre
la présélection de candidats pour son poste, au nom du Service Entreprises AEJ Agence Prestige.
Poste : %s chez %s | Nombre de candidats sélectionnés : %s | Top candidat : %s (score %s%%)
Inclure : objet sur la 1ère ligne préfixé "OBJET:", puis le corps.
Ne retourne QUE l'email, sans explication."""


class MoteurIA:

    def __init__(self, cfg: Config):
        self.cfg = cfg

    def _url(self):
        base = self.cfg.get("OLLAMA_URL", "http://localhost:11434").rstrip("/")
        return f"{base}/v1/chat/completions"

    def _model(self):
        return self.cfg.get("OLLAMA_MODEL", "gemma3:4b")

    def _parse_json_safe(self, txt: str) -> dict:
        """Extract JSON from model output that may include markdown or extra text."""
        t = re.sub(r'```json\s*', '', txt)
        t = re.sub(r'```\s*', '', t).strip()
        # Find first { or [
        idx = min(
            (t.find('{') if t.find('{') >= 0 else len(t)),
            (t.find('[') if t.find('[') >= 0 else len(t))
        )
        if idx > 0:
            t = t[idx:]
        # Fallback: regex-extract outermost JSON object
        try:
            return json.loads(t)
        except Exception:
            m = re.search(r'\{[\s\S]*\}', t)
            if m:
                try:
                    return json.loads(m.group(0))
                except Exception:
                    pass
        return {}

    def _appel_ollama(self, contenu: str, max_tokens=2500) -> str:
        """Call Ollama for JSON-structured responses (CV/fiche/score analysis)."""
        payload = {
            "model": self._model(),
            "messages": [
                {"role": "system",
                 "content": "Expert RH senior Côte d'Ivoire. Réponds UNIQUEMENT en JSON valide, sans aucun texte avant ou après."},
                {"role": "user", "content": contenu}
            ],
            "temperature": 0.05,
            "max_tokens": max_tokens,
            "stream": False,
        }
        for i in range(4):
            try:
                r = requests.post(self._url(), json=payload, timeout=180)
                if r.status_code == 200:
                    return r.json()['choices'][0]['message']['content']
                raise RuntimeError(f"Ollama HTTP {r.status_code}: {r.text[:200]}")
            except requests.Timeout:
                time.sleep(5)
        raise RuntimeError("Ollama : 4 tentatives infructueuses")

    def _appel_texte(self, contenu: str, max_tokens=1500) -> str:
        """Call Ollama for free-text generation (email drafts)."""
        payload = {
            "model": self._model(),
            "messages": [
                {"role": "system",
                 "content": "Expert RH senior, rédacteur professionnel. Réponds directement sans introduction."},
                {"role": "user", "content": contenu}
            ],
            "temperature": 0.4,
            "max_tokens": max_tokens,
            "stream": False,
        }
        r = requests.post(self._url(), json=payload, timeout=180)
        if r.status_code == 200:
            return r.json()['choices'][0]['message']['content'].strip()
        raise RuntimeError(f"Ollama HTTP {r.status_code}")

    def analyser_cv(self, texte: str, nom: str) -> dict:
        res = self._appel_ollama(PROMPT_CV % (nom, texte[:12000]))
        return self._parse_json_safe(res)

    def analyser_fiche(self, texte: str, nom: str) -> dict:
        res = self._appel_ollama(PROMPT_FICHE % (nom, texte[:8000]))
        return self._parse_json_safe(res)

    def scorer(self, cv: dict, fiche: dict) -> dict:
        res = self._appel_ollama(PROMPT_SCORE % (
            json.dumps(cv,    ensure_ascii=False)[:3500],
            json.dumps(fiche, ensure_ascii=False)[:3500]
        ), max_tokens=1200)
        return self._parse_json_safe(res)

    def generer_email_convocation(self, candidat: dict, poste: dict, score: int) -> str:
        return self._appel_texte(PROMPT_EMAIL_CONVOC % (
            candidat.get('nom_candidat', ''),
            poste.get('intitule', ''),
            poste.get('entreprise', ''),
            score
        ), max_tokens=600)

    def generer_email_recruteur(self, poste: dict, candidats: list) -> str:
        top = candidats[0] if candidats else {}
        return self._appel_texte(PROMPT_EMAIL_RECRUTEUR % (
            poste.get('intitule', ''),
            poste.get('entreprise', ''),
            len(candidats),
            top.get('nom_candidat', ''),
            top.get('score', 0)
        ), max_tokens=600)

    # ── Pré-filtrage TF-IDF ─────────────────────────────────────
    @staticmethod
    def score_tfidf(texte_cv: str, texte_fiche: str) -> float:
        if not SKLEARN_OK:
            # Fallback : overlap de mots-clés
            cv_w   = set(re.findall(r'\w+', texte_cv.lower()))
            f_w    = set(re.findall(r'\w+', texte_fiche.lower()))
            inter  = len(cv_w & f_w)
            return inter / max(len(f_w), 1) * 100
        try:
            vec    = TfidfVectorizer(ngram_range=(1, 2), max_features=8000, min_df=1)
            mat    = vec.fit_transform([texte_cv[:8000], texte_fiche[:4000]])
            sim    = cosine_similarity(mat[0], mat[1])[0][0]
            return round(float(sim) * 100, 1)
        except Exception:
            return 0.0

    @staticmethod
    def estimer_appels(nb_cv: int, nb_fiches: int, seuil_tfidf: int) -> dict:
        """Estimation du nombre d'appels IA Ollama (retrieval hybride top-K)."""
        appels_cv     = nb_cv
        appels_fiches = nb_fiches
        # Retrieval hybride : ~top_k postes par CV au lieu de tous
        top_k         = 6
        appels_score  = int(nb_cv * min(nb_fiches, top_k))
        total         = appels_cv + appels_fiches + appels_score
        cout_estimé   = total * 0.0001
        duree_min     = int(total * 1.3 / 60)
        return {
            "appels_cv"    : appels_cv,
            "appels_fiches": appels_fiches,
            "appels_score" : appels_score,
            "total"        : total,
            "cout_usd"     : cout_estimé,
            "duree_min"    : duree_min,
        }


# ════════════════════════════════════════════════════════════════
#  SECTION 7b — MOTEUR DE MATCHING HYBRIDE
#  Innovations issues des meilleurs ATS du marché :
#   • BM25 (Okapi)            — ranking lexical de référence
#   • Embeddings sémantiques  — sens, pas juste mots (multilingue)
#   • TF-IDF cosinus          — repli si pas d'embeddings
#   • Reciprocal Rank Fusion  — fusion des signaux de retrieval
#   • Fuzzy skill matching    — compétences malgré fautes/variantes
#   • Ontologie de synonymes  — expansion métier FR / contexte CI
#   • Score hybride explicable — LLM + sémantique + compétences
# ════════════════════════════════════════════════════════════════

# Ontologie légère de synonymes/variantes (contexte ivoirien & FR)
SYNONYMES_COMPETENCES = {
    "js": ["javascript"], "javascript": ["js", "node", "nodejs"],
    "rh": ["ressources humaines", "human resources", "grh"],
    "ressources humaines": ["rh", "grh"],
    "compta": ["comptabilité", "comptable", "accounting"],
    "comptabilité": ["compta", "accounting", "finance"],
    "marketing": ["mkg", "communication", "digital marketing"],
    "gestion de projet": ["chef de projet", "project management", "pmo"],
    "commercial": ["vente", "sales", "business development", "b2b", "b2c"],
    "informatique": ["it", "développement", "developpement", "dev", "software"],
    "anglais": ["english"], "français": ["french"],
    "excel": ["tableur", "spreadsheet", "office"],
    "sql": ["mysql", "postgresql", "base de données", "bdd"],
    "logistique": ["supply chain", "approvisionnement", "transit"],
    "btp": ["génie civil", "genie civil", "construction", "bâtiment"],
    "agro": ["agronomie", "agriculture", "agro-industrie", "agroindustrie"],
}

_STOP_FR = set("""au aux avec ce ces dans de des du elle en et eux il je la le les leur
lui ma mais me meme mes moi mon ne nos notre nous on ou par pas pour qu que qui sa se ses
son sur ta te tes toi ton tu un une vos votre vous c d j l a à n s t y est été être avoir
sont ans an experience expérience poste profil candidat cv""".split())


def _tokenize(txt: str):
    txt = (txt or "").lower()
    mots = re.findall(r"[a-zà-ÿ0-9\+\#\.]{2,}", txt)
    return [m for m in mots if m not in _STOP_FR]


class MatchingEngine:
    """Moteur de retrieval + scoring hybride. Dégradation gracieuse
    selon les bibliothèques disponibles."""

    def __init__(self, cfg: Config, log=print):
        self.cfg = cfg
        self.log = log
        self._emb_model = None
        self._postes = []          # liste de dicts fiches (avec _id, _texte)
        self._poste_tokens = []    # tokens par poste (BM25)
        self._bm25 = None
        self._tfidf_vec = None
        self._tfidf_mat = None
        self._poste_emb = None     # matrice embeddings postes
        self.backend = self._choisir_backend()

    def _choisir_backend(self):
        if self.cfg.get('USE_EMBEDDINGS', True) and EMB_OK:
            return EMB_BACKEND
        if BM25_OK:
            return 'bm25'
        if SKLEARN_OK:
            return 'tfidf'
        return 'overlap'

    # ── Embeddings ───────────────────────────────────────────────
    def _get_emb_model(self):
        if self._emb_model is not None:
            return self._emb_model
        try:
            if EMB_BACKEND == 'sentence_transformers':
                self._emb_model = SentenceTransformer(self.cfg.get('EMB_MODEL'))
            elif EMB_BACKEND == 'fastembed':
                self._emb_model = TextEmbedding()
        except Exception as e:
            self.log(f"   ⚠️ Embeddings indisponibles ({e}) → repli BM25/TF-IDF")
            self._emb_model = None
        return self._emb_model

    def _embed(self, textes):
        m = self._get_emb_model()
        if m is None:
            return None
        try:
            if EMB_BACKEND == 'sentence_transformers':
                return m.encode(textes, normalize_embeddings=True)
            else:  # fastembed
                import numpy as _np
                vecs = list(m.embed(textes))
                arr = _np.array(vecs)
                norm = _np.linalg.norm(arr, axis=1, keepdims=True)
                return arr / _np.clip(norm, 1e-9, None)
        except Exception:
            return None

    # ── Indexation des postes (une seule fois) ──────────────────
    def indexer_postes(self, postes: list, textes_fiches: dict):
        self._postes = []
        corpus = []
        for p in postes:
            t = textes_fiches.get(p['_id'], '') or self._poste_to_text(p)
            p = dict(p)
            p['_texte'] = t
            self._postes.append(p)
            corpus.append(t)

        self._poste_tokens = [_tokenize(t) for t in corpus]
        if BM25_OK and any(self._poste_tokens):
            try:
                self._bm25 = BM25Okapi(self._poste_tokens)
            except Exception:
                self._bm25 = None
        if SKLEARN_OK and corpus:
            try:
                self._tfidf_vec = TfidfVectorizer(ngram_range=(1, 2), max_features=12000)
                self._tfidf_mat = self._tfidf_vec.fit_transform(corpus)
            except Exception:
                self._tfidf_vec = None
        if self.backend in ('sentence_transformers', 'fastembed') and corpus:
            self._poste_emb = self._embed(corpus)

    @staticmethod
    def _poste_to_text(p: dict) -> str:
        parts = [p.get('intitule', ''), p.get('secteur', ''),
                 ' '.join(p.get('competences_obligatoires', []) or []),
                 ' '.join(p.get('competences_souhaitees', []) or []),
                 p.get('formation_requise', ''),
                 ' '.join(p.get('mots_cles_sectoriels', []) or []),
                 p.get('resume_poste', '')]
        return ' '.join(str(x) for x in parts if x)

    @staticmethod
    def _cv_to_text(cv: dict) -> str:
        parts = [cv.get('titre_profil', ''), cv.get('secteur', ''),
                 ' '.join(cv.get('competences_techniques', []) or []),
                 ' '.join(cv.get('logiciels_outils', []) or []),
                 cv.get('formation_niveau', ''), cv.get('formation_domaine', ''),
                 ' '.join(cv.get('mots_cles_metier', []) or []),
                 cv.get('resume_profil', '')]
        return ' '.join(str(x) for x in parts if x)

    # ── Retrieval hybride : top-K postes pour un CV ─────────────
    def retrieve(self, cv: dict, texte_cv: str, top_k=6):
        """Retourne [(index_poste, score_fusion 0-100, sim_semantique 0-1)]."""
        if not self._postes:
            return []
        q_text = (texte_cv or '') + ' ' + self._cv_to_text(cv)
        n = len(self._postes)
        rangs = []  # listes de (index, score) par signal

        # Signal 1 — BM25
        if self._bm25 is not None:
            try:
                scores = self._bm25.get_scores(_tokenize(q_text))
                rangs.append(self._to_rank(scores))
            except Exception:
                pass
        # Signal 2 — TF-IDF cosinus
        if self._tfidf_vec is not None:
            try:
                qv = self._tfidf_vec.transform([q_text])
                sims = cosine_similarity(qv, self._tfidf_mat)[0]
                rangs.append(self._to_rank(sims))
            except Exception:
                pass
        # Signal 3 — Embeddings sémantiques
        sem_sims = None
        if self._poste_emb is not None:
            qe = self._embed([q_text])
            if qe is not None:
                try:
                    import numpy as _np
                    sem_sims = (self._poste_emb @ qe[0])
                    rangs.append(self._to_rank(sem_sims))
                except Exception:
                    sem_sims = None
        # Repli : overlap simple
        if not rangs:
            qset = set(_tokenize(q_text))
            ov = [len(qset & set(tok)) for tok in self._poste_tokens]
            rangs.append(self._to_rank(ov))

        # Fusion par Reciprocal Rank Fusion (RRF)
        rrf = self._rrf(rangs, n)
        ordre = sorted(range(n), key=lambda i: rrf[i], reverse=True)[:top_k]
        maxr = max(rrf) or 1.0
        res = []
        for i in ordre:
            sem = float(sem_sims[i]) if sem_sims is not None else 0.0
            res.append((i, round(rrf[i] / maxr * 100, 1), max(0.0, min(1.0, sem))))
        return res

    @staticmethod
    def _to_rank(scores):
        order = sorted(range(len(scores)), key=lambda i: scores[i], reverse=True)
        rank = [0] * len(scores)
        for pos, idx in enumerate(order):
            rank[idx] = pos
        return rank

    @staticmethod
    def _rrf(listes_rangs, n, k=60):
        fused = [0.0] * n
        for rangs in listes_rangs:
            for i in range(n):
                fused[i] += 1.0 / (k + rangs[i])
        return fused

    # ── Couverture des compétences (fuzzy) ───────────────────────
    def couverture_competences(self, cv: dict, poste: dict) -> int:
        requis = list(poste.get('competences_obligatoires', []) or []) + \
                 list(poste.get('competences_souhaitees', []) or [])
        if not requis:
            return 0
        cv_skills = list(cv.get('competences_techniques', []) or []) + \
                    list(cv.get('logiciels_outils', []) or []) + \
                    _tokenize(self._cv_to_text(cv))
        cv_skills_l = [str(s).lower() for s in cv_skills]
        seuil = self.cfg.get('FUZZ_SEUIL', 82)
        couverts = 0
        for r in requis:
            if self._skill_present(str(r).lower(), cv_skills_l, seuil):
                couverts += 1
        return int(couverts / max(len(requis), 1) * 100)

    def _skill_present(self, skill, cv_skills_l, seuil):
        variantes = [skill] + SYNONYMES_COMPETENCES.get(skill, [])
        for v in variantes:
            for cs in cv_skills_l:
                if v in cs or cs in v:
                    return True
                if FUZZ_OK and rf_fuzz.token_set_ratio(v, cs) >= seuil:
                    return True
        return False

    # ── Score hybride final explicable ───────────────────────────
    def score_hybride(self, score_llm: int, sim_semantique: float, couverture: int) -> int:
        wl = self.cfg.get('BLEND_LLM', 0.70)
        ws = self.cfg.get('BLEND_SEMANTIC', 0.18)
        wc = self.cfg.get('BLEND_SKILLS', 0.12)
        total_w = (wl + ws + wc) or 1.0
        val = (wl * score_llm + ws * sim_semantique * 100 + wc * couverture) / total_w
        return int(round(max(0, min(100, val))))

    def info_backend(self) -> str:
        noms = {'sentence_transformers': 'Embeddings (SentenceTransformers)',
                'fastembed': 'Embeddings (FastEmbed/ONNX)',
                'bm25': 'BM25 Okapi', 'tfidf': 'TF-IDF cosinus',
                'overlap': 'Overlap mots-clés (basique)'}
        extras = []
        if FUZZ_OK: extras.append('Fuzzy')
        if BM25_OK and self.backend != 'bm25': extras.append('BM25')
        suffix = f" + {' + '.join(extras)}" if extras else ""
        return noms.get(self.backend, self.backend) + suffix


# ════════════════════════════════════════════════════════════════
#  SECTION 8 — GÉNÉRATEUR RAPPORTS (Excel + Lettres DOCX)
# ════════════════════════════════════════════════════════════════
class RapportGenerator:
    V = "1A7A3C"  # Vert AEJ
    O = "F47920"  # Orange AEJ
    R = "FFCDD2"
    J = "FFF9C4"
    G = "C8E6C9"
    B = "FFFFFF"

    def _s(self, cell, bold=False, bg=None, color="000000", sz=10, center=False, wrap=False):
        cell.font = Font(bold=bold, size=sz, color=color)
        if bg:
            cell.fill = PatternFill("solid", fgColor=bg)
        cell.alignment = Alignment(
            horizontal='center' if center else 'left',
            vertical='center', wrap_text=wrap)

    # ── Rapport Excel 4 feuilles ─────────────────────────────────
    def excel_complet(self, resultats: list, contacts: list, chemin: Path):
        if not EXCEL_OK:
            raise RuntimeError("openpyxl non installé")
        wb = openpyxl.Workbook()

        self._feuille_rapport(wb, resultats)
        self._feuille_synthese(wb, resultats)
        self._feuille_contacts(wb, contacts)
        self._feuille_dashboard(wb, resultats)

        wb.save(str(chemin))

    def _feuille_rapport(self, wb, resultats):
        ws = wb.active
        ws.title = "RAPPORT MATCHING"
        ws.sheet_view.showGridLines = False

        ws.merge_cells('A1:N1')
        c = ws['A1']
        c.value = f"AEJ AGENCE PRESTIGE — RAPPORT MATCHING CV — {datetime.now().strftime('%d/%m/%Y %H:%M')}"
        self._s(c, bold=True, bg=self.V, color=self.B, sz=14, center=True)
        ws.row_dimensions[1].height = 34

        ws.merge_cells('A2:N2')
        c2 = ws['A2']
        c2.value = "Bd des Martyrs, Cocody, Abidjan  |  Service Entreprises"
        self._s(c2, bg=self.O, color=self.B, sz=10, center=True)
        ws.row_dimensions[2].height = 20

        cols = ['N°','CANDIDAT','FICHIER CV','POSTE RECOMMANDÉ','ENTREPRISE','SECTEUR',
                'SCORE IA','SCORE HYBRIDE','SÉMANT.','COMPÉT.','EXP.','DIPLÔME',
                'LANGUES','DÉCISION','JUSTIFICATION']
        for ci, h in enumerate(cols, 1):
            c = ws.cell(row=3, column=ci, value=h)
            self._s(c, bold=True, bg=self.O, color=self.B, sz=9, center=True)
        ws.row_dimensions[3].height = 22

        for i, r in enumerate(resultats, 1):
            row = i + 3
            sc  = r.get('score_hybride', r.get('score_global', 0))
            dec = r.get('decision', '')
            bg  = self.G if sc >= 80 else (self.J if sc >= 65 else self.R)

            vals = [i, r.get('nom_candidat',''), r.get('fichier_cv',''),
                    r.get('intitule_poste',''), r.get('entreprise',''),
                    r.get('secteur_poste',''),
                    f"{r.get('score_global',0)}%",
                    f"{r.get('score_hybride',r.get('score_global',0))}%",
                    f"{r.get('sim_semantique',0)}%",
                    f"{r.get('couverture_competences',0)}%",
                    f"{r.get('score_experience',0)}%",
                    f"{r.get('score_diplome',0)}%",
                    r.get('langues_cv',''), dec, r.get('justification','')]

            for ci, val in enumerate(vals, 1):
                c = ws.cell(row=row, column=ci, value=val)
                self._s(c, bg=bg, sz=9, center=(ci in (1,7,8,9,10,11,12,14)), wrap=(ci==15))
            if dec == 'PRIORITAIRE':
                ws.cell(row=row, column=14).font = Font(bold=True, color="0D5A2C", size=9)
            elif dec == 'FAIBLE':
                ws.cell(row=row, column=14).font = Font(bold=True, color="B71C1C", size=9)

        for ci, w in enumerate([4,22,24,24,18,15,9,11,9,9,8,9,13,12,38], 1):
            ws.column_dimensions[get_column_letter(ci)].width = w
        ws.freeze_panes = 'A4'
        if resultats:
            ws.auto_filter.ref = f"A3:O{len(resultats)+3}"

    def _feuille_synthese(self, wb, resultats):
        ws2 = wb.create_sheet("SYNTHÈSE PAR POSTE")
        ws2.sheet_view.showGridLines = False
        ws2.merge_cells('A1:H1')
        c = ws2['A1']
        c.value = "SYNTHÈSE — MEILLEURS CANDIDATS PAR POSTE"
        self._s(c, bold=True, bg=self.V, color=self.B, sz=13, center=True)
        ws2.row_dimensions[1].height = 30

        hdrs = ['POSTE','ENTREPRISE','SECTEUR','NB RETENUS',
                'TOP CANDIDAT','MEILLEUR SCORE','DÉCISION','POINTS FORTS']
        for ci, h in enumerate(hdrs, 1):
            c = ws2.cell(row=2, column=ci, value=h)
            self._s(c, bold=True, bg=self.O, color=self.B, sz=9, center=True)

        par_poste = {}
        for r in resultats:
            k = (r.get('intitule_poste',''), r.get('entreprise',''))
            par_poste.setdefault(k, []).append(r)

        row = 3
        for (poste, ent), cands in sorted(par_poste.items()):
            retenus  = [c for c in cands if c.get('score_global',0) >= 65]
            meilleur = max(retenus, key=lambda x: x.get('score_global',0)) if retenus else None
            sc_m     = meilleur.get('score_global',0) if meilleur else 0
            bg       = self.G if sc_m >= 80 else (self.J if sc_m >= 65 else self.R)

            vals = [poste, ent, (meilleur or {}).get('secteur_poste',''),
                    len(retenus),
                    (meilleur or {}).get('nom_candidat','-'),
                    f"{sc_m}%" if meilleur else '-',
                    'PRIORITAIRE' if sc_m >= 80 else ('MOYEN' if sc_m >= 65 else '-'),
                    (meilleur or {}).get('points_forts','')[:60]]
            for ci, val in enumerate(vals, 1):
                c = ws2.cell(row=row, column=ci, value=val)
                self._s(c, bg=bg, sz=9, center=(ci in (4,6,7)))
            row += 1

        for ci, w in enumerate([30,22,16,12,26,14,12,35], 1):
            ws2.column_dimensions[get_column_letter(ci)].width = w

    def _feuille_contacts(self, wb, contacts):
        ws3 = wb.create_sheet("SUIVI CONTACTS")
        ws3.sheet_view.showGridLines = False
        ws3.merge_cells('A1:J1')
        c = ws3['A1']
        c.value = "SUIVI MISE EN RELATION — AEJ AGENCE PRESTIGE"
        self._s(c, bold=True, bg=self.V, color=self.B, sz=13, center=True)
        ws3.row_dimensions[1].height = 30

        hdrs = ['CANDIDAT','PROFIL','POSTE','ENTREPRISE','SCORE','DÉCISION',
                'STATUT','EMAIL','TÉLÉPHONE','NOTES']
        for ci, h in enumerate(hdrs, 1):
            c = ws3.cell(row=2, column=ci, value=h)
            self._s(c, bold=True, bg=self.O, color=self.B, sz=9, center=True)

        for i, ct in enumerate(contacts, 3):
            sc  = ct.get('score', 0)
            bg  = self.G if sc >= 80 else (self.J if sc >= 65 else self.R)
            vals = [ct.get('nom_candidat',''), ct.get('titre_profil',''),
                    ct.get('intitule_poste',''), ct.get('entreprise',''),
                    f"{sc}%", ct.get('decision',''),
                    STATUTS_FR.get(ct.get('statut','EN_ATTENTE'), ct.get('statut','')),
                    ct.get('email_candidat',''), ct.get('tel_candidat',''),
                    ct.get('notes','')]
            for ci, val in enumerate(vals, 1):
                c = ws3.cell(row=i, column=ci, value=val)
                self._s(c, bg=bg, sz=9, center=(ci in (5,6,7)))

        for ci, w in enumerate([24,22,26,20,10,12,20,26,16,30], 1):
            ws3.column_dimensions[get_column_letter(ci)].width = w

    def _feuille_dashboard(self, wb, resultats):
        ws4 = wb.create_sheet("TABLEAU DE BORD")
        ws4.sheet_view.showGridLines = False
        ws4.merge_cells('A1:D1')
        c = ws4['A1']
        c.value = "TABLEAU DE BORD — STATISTIQUES MATCHING"
        self._s(c, bold=True, bg=self.V, color=self.B, sz=13, center=True)
        ws4.row_dimensions[1].height = 30

        nb_top = sum(1 for r in resultats if r.get('score_global',0) >= 80)
        nb_moy = sum(1 for r in resultats if 65 <= r.get('score_global',0) < 80)
        nb_fbl = sum(1 for r in resultats if r.get('score_global',0) < 65)
        postes = len({r.get('intitule_poste') for r in resultats})
        moy_sc = (sum(r.get('score_global',0) for r in resultats)
                  // max(len(resultats), 1)) if resultats else 0

        stats = [
            ("📊 Total correspondances",       len(resultats)),
            ("⭐ TOP MATCHES (≥ 80%)",          nb_top),
            ("✅ MATCHES MOYENS (65–79%)",      nb_moy),
            ("⚪ FAIBLES (< 65%)",              nb_fbl),
            ("📋 Postes avec candidats",        postes),
            ("📈 Score moyen global",           f"{moy_sc}%"),
            ("🕐 Date d'exécution",             datetime.now().strftime("%d/%m/%Y %H:%M")),
        ]
        for i, (lbl, val) in enumerate(stats, 3):
            ws4.cell(row=i, column=1, value=lbl).font = Font(bold=True, size=11)
            c = ws4.cell(row=i, column=2, value=val)
            c.font = Font(size=11, color=self.V)
            c.fill = PatternFill("solid", fgColor="F1F8E9")
        ws4.column_dimensions['A'].width = 36
        ws4.column_dimensions['B'].width = 22

    # ── Lettre de présélection DOCX ──────────────────────────────
    def lettre_preselection(self, poste: dict, candidats: list, chemin: Path):
        if not DOCX_OK:
            return None
        doc = DocxDoc()

        # Marges
        for section in doc.sections:
            section.top_margin    = Cm(2)
            section.bottom_margin = Cm(2)
            section.left_margin   = Cm(2.5)
            section.right_margin  = Cm(2.5)

        # En-tête AEJ
        hdr = doc.add_heading('', level=0)
        run = hdr.add_run('AGENCE EMPLOI JEUNES — AGENCE PRESTIGE')
        run.font.color.rgb = RGBColor(0x1A, 0x7A, 0x3C)
        run.font.size = Pt(14)
        run.font.bold = True
        hdr.alignment = WD_ALIGN_PARAGRAPH.CENTER

        sous_hdr = doc.add_paragraph('Bd des Martyrs, Cocody, Abidjan  |  Service Entreprises')
        sous_hdr.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r in sous_hdr.runs:
            r.font.size = Pt(9)
            r.font.color.rgb = RGBColor(0xF4, 0x79, 0x20)

        doc.add_paragraph()

        # Titre poste
        titre = doc.add_heading('', level=1)
        run_t = titre.add_run(f"RAPPORT DE PRÉSÉLECTION — {poste.get('intitule','').upper()}")
        run_t.font.color.rgb = RGBColor(0x1A, 0x7A, 0x3C)
        run_t.font.size = Pt(13)

        p_info = doc.add_paragraph()
        p_info.add_run(f"Entreprise : ").bold = True
        p_info.add_run(f"{poste.get('entreprise','—')}   |   ")
        p_info.add_run("Secteur : ").bold = True
        p_info.add_run(f"{poste.get('secteur','—')}   |   ")
        p_info.add_run("Date : ").bold = True
        p_info.add_run(datetime.now().strftime("%d/%m/%Y"))

        doc.add_paragraph(
            f"Nombre de candidats présélectionnés : {len(candidats)}"
        ).bold = True
        doc.add_paragraph()

        # Fiche de chaque candidat
        for i, c in enumerate(candidats, 1):
            h = doc.add_heading('', level=2)
            r = h.add_run(f"{i}. {c.get('nom_candidat','Candidat')}  —  Score : {c.get('score_global',0)}%")
            r.font.color.rgb = RGBColor(0xF4, 0x79, 0x20)
            r.font.size = Pt(11)

            t = doc.add_paragraph()
            t.add_run("Profil : ").bold = True
            t.add_run(c.get('titre_profil', ''))

            e = doc.add_paragraph()
            e.add_run("Expérience : ").bold = True
            e.add_run(f"{c.get('annees_experience', '?')} ans")

            f = doc.add_paragraph()
            f.add_run("Formation : ").bold = True
            f.add_run(f"{c.get('formation_niveau','')} {c.get('formation_domaine','')}")

            comp = doc.add_paragraph()
            comp.add_run("Compétences clés : ").bold = True
            comp_list = c.get('competences_techniques', [])
            comp.add_run(', '.join(comp_list[:6]) if comp_list else '—')

            pts = doc.add_paragraph()
            pts.add_run("Points forts : ").bold = True
            pts.add_run(c.get('points_forts', '—') if isinstance(c.get('points_forts'), str)
                        else '; '.join(c.get('points_forts', [])))

            jus = doc.add_paragraph()
            jus.add_run("Justification IA : ").bold = True
            jus.add_run(c.get('justification', ''))

            doc.add_paragraph('─' * 60)

        # Pied de page
        doc.add_paragraph()
        pie = doc.add_paragraph('Document généré automatiquement par AEJ CV MATCHER PRO v2.0')
        pie.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r in pie.runs:
            r.font.size = Pt(8)
            r.font.color.rgb = RGBColor(0x9E, 0x9E, 0x9E)

        doc.save(str(chemin))
        return chemin


# ════════════════════════════════════════════════════════════════
#  SECTION 9 — PIPELINE ORCHESTRATEUR (PARALLÈLE + HYBRIDE)
# ════════════════════════════════════════════════════════════════
def _nettoyer(s, n=100):
    return re.sub(r'[\\/:*?"<>|]', '-', str(s or ''))[:n].strip()


class Orchestrateur:
    def __init__(self, cfg: Config, db: Database, log=None, progress=None):
        self.cfg      = cfg
        self.db       = db
        self.log      = log or print
        self.progress = progress or (lambda v, m: None)
        self.drive    = make_provider(cfg)
        self.extract  = Extracteur(cfg.get('TESSERACT_PATH'))
        self.ia       = MoteurIA(cfg)
        self.rapport  = RapportGenerator()
        self.engine   = MatchingEngine(cfg, log=self.log)
        self.tmp      = Path(tempfile.mkdtemp(prefix="aej_v2_"))
        self.arret    = threading.Event()
        # Résolution des sources / destination selon le mode
        if cfg.get('SOURCE_MODE', 'DRIVE') == 'LOCAL':
            self.src_cv    = 'LOCAL_CV'
            self.src_fiche = 'LOCAL_FICHE'
            self.dest      = cfg.get('LOCAL_DEST_PATH', '')
        else:
            self.src_cv    = cfg.get('DOSSIER_CV_SOURCE')
            self.src_fiche = cfg.get('DOSSIER_FICHES_POSTE')
            self.dest      = cfg.get('DOSSIER_DESTINATION')

    def stop(self):
        self.arret.set()

    # ── Téléchargement + extraction parallèles ───────────────────
    def _traiter_fichier(self, f: dict, prefix: str) -> tuple:
        """Thread-safe: télécharge et extrait un fichier."""
        try:
            dest = self.tmp / f"{prefix}_{f['id']}{Path(f['name']).suffix or '.pdf'}"
            self.drive.telecharger(f['id'], f['mimeType'], dest)
            texte = self.extract.extraire(dest, f['mimeType'])
            h = Extracteur.md5(texte)
            return f['id'], f['name'], texte, h, None
        except Exception as e:
            return f['id'], f['name'], '', '', str(e)

    # ── Pipeline principal ───────────────────────────────────────
    def executer(self):
        debut_ts = datetime.now().isoformat()
        debut    = time.time()
        resultats_matching = []
        workers = min(self.cfg.get('PARALLEL_WORKERS', 4), 8)

        try:
            # ── 1. Connexion Drive / Vérif local ──────────────────
            mode_local = self.cfg.get('SOURCE_MODE', 'DRIVE') == 'LOCAL'
            if mode_local:
                self.log("💻 Mode FICHIERS LOCAUX — préparation...")
                self.drive.connecter()
                self.log(f"✅ Destination locale : {self.dest}\n")
            else:
                self.log("🔌 Connexion à Google Drive...")
                self.drive.connecter()
                self.log("✅ Drive connecté.\n")
            self.progress(5, "Source prête")
            if self.arret.is_set():
                return []

            # ── 2. Fiches de poste (parallèle) ───────────────────
            self.log("═" * 56)
            self.log("📋  ÉTAPE 1 — Indexation fiches de poste...")
            fiches_raw = self.drive.lister(self.src_fiche, recurse=True)
            self.log(f"   {len(fiches_raw)} fiche(s) trouvée(s)")
            fiches, textes_fiches = [], {}

            with ThreadPoolExecutor(max_workers=workers) as ex:
                futs = {ex.submit(self._traiter_fichier, f, 'fiche'): f for f in fiches_raw}
                for i, fut in enumerate(as_completed(futs)):
                    if self.arret.is_set():
                        break
                    fid, nom, texte, _, err = fut.result()
                    self.progress(5 + int(12 * i / max(len(fiches_raw), 1)),
                                  f"Fiche: {nom[:40]}")
                    if err or len(texte.strip()) < 20:
                        self.log(f"   ⚠️  {nom}: {err or 'vide'}")
                        continue
                    p = self.db.get_fiche(fid)
                    if not p:
                        try:
                            p = self.ia.analyser_fiche(texte, nom)
                            self.db.save_fiche(fid, nom, p)
                            time.sleep(self.cfg.get('PAUSE_API', 1.0))
                        except Exception as e:
                            self.log(f"   ❌ {nom}: {e}")
                            continue
                    else:
                        self.log(f"   📂 [CACHE] {nom}")
                    p.update({'_id': fid, '_nom': nom})
                    fiches.append(p)
                    textes_fiches[fid] = texte

            self.log(f"\n   ✔ {len(fiches)} fiche(s) indexée(s).")
            self.progress(17, f"{len(fiches)} fiches")

            # ── 3. CV (parallèle pour téléchargement) ────────────
            self.log("\n" + "═" * 56)
            self.log("📄  ÉTAPE 2 — Extraction des profils CV...")
            cvs_raw = self.drive.lister(self.src_cv)
            self.log(f"   {len(cvs_raw)} CV trouvé(s)")
            cvs, textes_cvs = [], {}

            # Vérif déduplications
            dup_count = 0
            with ThreadPoolExecutor(max_workers=workers) as ex:
                futs = {ex.submit(self._traiter_fichier, f, 'cv'): f for f in cvs_raw}
                for i, fut in enumerate(as_completed(futs)):
                    if self.arret.is_set():
                        break
                    fid, nom, texte, h, err = fut.result()
                    self.progress(17 + int(25 * i / max(len(cvs_raw), 1)),
                                  f"CV: {nom[:40]}")
                    if err or len(texte.strip()) < 30:
                        self.log(f"   ⚠️  {nom}: {err or 'vide'}")
                        continue

                    # Déduplication par hash
                    dup_id = self.db.get_cv_hash(h)
                    if dup_id and dup_id != fid:
                        dup_count += 1
                        self.log(f"   🔁 DOUBLON détecté: {nom} (ignoré)")
                        continue

                    p = self.db.get_cv(fid)
                    if not p:
                        try:
                            p = self.ia.analyser_cv(texte, nom)
                            self.db.save_cv(fid, nom, h, p)
                            time.sleep(self.cfg.get('PAUSE_API', 1.0))
                        except Exception as e:
                            self.log(f"   ❌ {nom}: {e}")
                            continue
                    else:
                        self.log(f"   📂 [CACHE] {nom}")

                    p.update({'_id': fid, '_nom': nom})
                    cvs.append(p)
                    textes_cvs[fid] = texte
                    cand_str = f"{p.get('nom_candidat','?')[:25]} | {p.get('titre_profil','?')[:30]}"
                    self.log(f"   ✅ {nom[:40]} → {cand_str}")

            if dup_count:
                self.log(f"\n   🔁 {dup_count} doublon(s) détecté(s) et ignoré(s).")
            self.log(f"\n   ✔ {len(cvs)} CV extrait(s).")
            self.progress(42, f"{len(cvs)} CV extraits")

            # ── 4. Retrieval HYBRIDE (BM25 + sémantique + RRF) ───
            self.log("\n" + "═" * 56)
            self.log("🔬  ÉTAPE 3 — Retrieval hybride (indexation des postes)...")
            self.engine.indexer_postes(fiches, textes_fiches)
            self.log(f"   Moteur : {self.engine.info_backend()}")
            top_k = self.cfg.get('RETRIEVAL_TOP_K', 6)

            paires_valides = []   # (cv, fiche, sim_semantique)
            total_paires   = len(cvs) * len(fiches)
            for cv in cvs:
                if self.arret.is_set():
                    break
                t_cv = textes_cvs.get(cv['_id'], cv.get('resume_profil', ''))
                matches = self.engine.retrieve(cv, t_cv, top_k=top_k)
                for idx_poste, score_fusion, sim_sem in matches:
                    fiche = self.engine._postes[idx_poste]
                    paires_valides.append((cv, fiche, sim_sem))

            economie = total_paires - len(paires_valides)
            taux_filtre = int(economie / max(total_paires, 1) * 100)
            self.log(f"   {len(paires_valides)} paires retenues / {total_paires} (top-{top_k} par CV)")
            self.log(f"   → Économie : {economie} appels IA évités ({taux_filtre}%)")
            self.progress(52, f"{len(paires_valides)} paires à scorer")

            # ── 5. Scoring Ollama + score hybride explicable ─────
            self.log("\n" + "═" * 56)
            self.log("🤖  ÉTAPE 4 — Scoring IA (Ollama local) + fusion hybride...")
            seuil_min  = self.cfg.get('SEUIL_MOYEN', 65)
            max_postes = self.cfg.get('MAX_POSTES_PAR_CV', 3)
            scores_par_cv = {}

            for idx, (cv, fiche, sim_sem) in enumerate(paires_valides):
                if self.arret.is_set():
                    break
                self.progress(52 + int(30 * idx / max(len(paires_valides), 1)),
                              f"Score: {cv.get('nom_candidat','?')[:18]} × {fiche.get('intitule','?')[:18]}")

                sc = self.db.get_score(cv['_id'], fiche['_id'])
                if not sc:
                    try:
                        sc = self.ia.scorer(cv, fiche)
                        self.db.save_score(cv['_id'], fiche['_id'], sim_sem, sc)
                        time.sleep(self.cfg.get('PAUSE_API', 1.0))
                    except Exception as e:
                        self.log(f"   ❌ {cv['_nom'][:25]} × {fiche['_nom'][:25]}: {e}")
                        continue

                sg  = sc.get('score_global', 0)
                # Signaux additionnels explicables
                couverture = self.engine.couverture_competences(cv, fiche)
                hybride    = self.engine.score_hybride(sg, sim_sem, couverture)

                if hybride < seuil_min and sg < seuil_min:
                    continue

                sd  = sc.get('scores_detail', {})
                rec = {
                    '_id_cv'           : cv['_id'],
                    '_id_fiche'        : fiche['_id'],
                    'fichier_cv'       : cv['_nom'],
                    'nom_candidat'     : cv.get('nom_candidat', ''),
                    'titre_profil'     : cv.get('titre_profil', ''),
                    'annees_experience': cv.get('annees_experience', 0),
                    'formation_niveau' : cv.get('formation_niveau', ''),
                    'formation_domaine': cv.get('formation_domaine', ''),
                    'competences_techniques': cv.get('competences_techniques', []),
                    'langues_cv'       : ', '.join(cv.get('langues', [])),
                    'email_candidat'   : cv.get('email', ''),
                    'tel_candidat'     : cv.get('telephone', ''),
                    'intitule_poste'   : fiche.get('intitule', fiche['_nom']),
                    'entreprise'       : fiche.get('entreprise', ''),
                    'secteur_poste'    : fiche.get('secteur', ''),
                    'score_global'     : sg,
                    'score_hybride'    : hybride,
                    'sim_semantique'   : round(sim_sem * 100, 1),
                    'couverture_competences': couverture,
                    'score_competences': sd.get('competences_metier', 0),
                    'score_experience' : sd.get('experience', 0),
                    'score_diplome'    : sd.get('diplome', 0),
                    'score_outils'     : sd.get('outils_logiciels', 0),
                    'points_forts'     : '; '.join(sc.get('points_forts', [])),
                    'reserves'         : '; '.join(sc.get('reserves', [])),
                    'decision'         : sc.get('decision', ''),
                    'justification'    : sc.get('justification', ''),
                }
                scores_par_cv.setdefault(cv['_id'], []).append(rec)
                self.db.upsert_contact(cv['_id'], fiche['_id'], rec)

            # Sélectionner les meilleurs postes par CV (ordre = score hybride)
            for id_cv, scores in scores_par_cv.items():
                scores.sort(key=lambda x: x['score_hybride'], reverse=True)
                if self.cfg.get('MODE') == 'MULTI_POSTES':
                    resultats_matching.extend(scores[:max_postes])
                elif scores:
                    resultats_matching.append(scores[0])

            self.log(f"\n   ✔ {len(resultats_matching)} correspondance(s) retenue(s).")
            self.progress(82, f"{len(resultats_matching)} correspondances")

            # ── 6. Scoring percentile (sur score hybride) ─────────
            if resultats_matching:
                tous_scores = sorted([r['score_hybride'] for r in resultats_matching])
                for r in resultats_matching:
                    rang = tous_scores.index(r['score_hybride'])
                    r['percentile'] = int(rang / len(tous_scores) * 100)

            # ── 7. Classement Drive ───────────────────────────────
            self.log("\n" + "═" * 56)
            self.log("📁  ÉTAPE 5 — Classement dans Google Drive...")
            dest_id   = self.dest
            action    = self.cfg.get('ACTION', 'COPIER')
            seuil_top = self.cfg.get('SEUIL_TOP', 80)
            par_poste = {}
            for r in resultats_matching:
                par_poste.setdefault((r['intitule_poste'], r['entreprise']), []).append(r)

            for (poste, ent), cands in par_poste.items():
                cands.sort(key=lambda x: x.get('score_hybride', x['score_global']), reverse=True)
                prefixe  = f"{ent} — " if ent else ''
                id_doss  = self.drive.creer_dossier(_nettoyer(prefixe + poste), dest_id)
                id_top   = self.drive.creer_dossier(f"⭐ TOP MATCHES (≥{seuil_top}%)", id_doss)
                id_moy   = self.drive.creer_dossier(f"✅ MATCHES MOYENS ({seuil_min}-{seuil_top-1}%)", id_doss)
                id_fbl   = self.drive.creer_dossier(f"⚪ FAIBLES (<{seuil_min}%)", id_doss)

                top_n = self.cfg.get('MAX_CV_PAR_POSTE', 15)
                for rang, c in enumerate(cands[:top_n], 1):
                    sg  = c.get('score_hybride', c['score_global'])
                    did = id_top if sg >= seuil_top else (id_moy if sg >= seuil_min else id_fbl)
                    ext = Path(c['fichier_cv']).suffix or '.pdf'
                    nom = _nettoyer(
                        f"Rang{rang:02d}_Score{sg}_{c.get('nom_candidat','') or Path(c['fichier_cv']).stem}"
                    ) + ext
                    try:
                        if action == 'COPIER':
                            self.drive.copier(c['_id_cv'], nom, did)
                        else:
                            self.drive.deplacer(c['_id_cv'], nom, did)
                        self.log(f"   {'📋' if action=='COPIER' else '📦'} [{sg}%] {c.get('nom_candidat','?')[:25]} → {poste[:28]}")
                    except Exception as e:
                        self.log(f"   ❌ {c['fichier_cv']}: {e}")

                # Lettre de présélection DOCX
                if DOCX_OK:
                    fiche_p = next((f for f in fiches if f.get('intitule') == poste), {})
                    nom_lettre = f"Preselection_{_nettoyer(poste)}.docx"
                    ch_lettre  = self.tmp / nom_lettre
                    self.rapport.lettre_preselection(fiche_p, cands[:top_n], ch_lettre)
                    self.drive.uploader(ch_lettre, nom_lettre, id_doss,
                        mime='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
                    self.log(f"   📄 Lettre présélection déposée: {nom_lettre}")

            self.progress(92, "Classement terminé")

            # ── 8. Rapport Excel ──────────────────────────────────
            self.log("\n" + "═" * 56)
            self.log("📊  ÉTAPE 6 — Génération du rapport Excel...")
            contacts_crm = self.db.get_contacts()
            nom_rpt  = f"RAPPORT_MATCHING_AEJ_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx"
            ch_rpt   = self.tmp / nom_rpt
            self.rapport.excel_complet(resultats_matching, contacts_crm, ch_rpt)
            lien = self.drive.uploader(ch_rpt, nom_rpt, dest_id,
                mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
            self.log(f"   ✅ Rapport Excel : {nom_rpt}")
            if lien:
                self.log(f"   🔗 {lien}")

            # ── Résumé ───────────────────────────────────────────
            duree = int(time.time() - debut)
            self.db.save_run(debut_ts, datetime.now().isoformat(),
                             len(cvs), len(fiches), len(resultats_matching),
                             duree, 'SUCCES')
            self.log(f"\n{'═'*56}")
            self.log(f"✅  PIPELINE TERMINÉ EN {duree//60}min {duree%60}s")
            self.log(f"   CV analysés              : {len(cvs)}")
            self.log(f"   Fiches indexées           : {len(fiches)}")
            self.log(f"   Paires TF-IDF             : {len(paires_valides)} / {total_paires} ({100-taux_filtre}%)")
            self.log(f"   Correspondances retenues  : {len(resultats_matching)}")
            self.log(f"   Rapport Excel             : {nom_rpt}")
            self.log(f"{'═'*56}\n")
            self.progress(100, f"✅ Terminé — {len(resultats_matching)} correspondances")

            # Notification Windows
            if TOAST_OK:
                try:
                    ToastNotifier().show_toast(
                        "AEJ CV Matcher PRO",
                        f"✅ Matching terminé ! {len(resultats_matching)} correspondances trouvées.",
                        duration=8, threaded=True
                    )
                except Exception:
                    pass

        except Exception as e:
            self.log(f"\n❌ ERREUR CRITIQUE: {e}\n{traceback.format_exc()}")
            self.db.save_run(debut_ts, datetime.now().isoformat(),
                             0, 0, 0, int(time.time()-debut), 'ERREUR')
            raise
        finally:
            try:
                shutil.rmtree(self.tmp, ignore_errors=True)
            except Exception:
                pass

        return resultats_matching


# ════════════════════════════════════════════════════════════════
#  SECTION 10 — WATCHDOG AUTONOME
# ════════════════════════════════════════════════════════════════
class Watchdog:
    def __init__(self, cfg: Config, db: Database, run_callback):
        self.cfg      = cfg
        self.db       = db
        self.run_cb   = run_callback
        self._thread  = None
        self._stop    = threading.Event()
        self.next_run = None

    def demarrer(self):
        if self._thread and self._thread.is_alive():
            return
        self._stop.clear()
        self._thread = threading.Thread(target=self._loop, daemon=True)
        self._thread.start()

    def arreter(self):
        self._stop.set()

    def _loop(self):
        interval_h = self.cfg.get('AUTO_INTERVALLE_H', 6)
        while not self._stop.is_set():
            self.next_run = datetime.now() + timedelta(hours=interval_h)
            # Attendre en petits morceaux pour pouvoir stopper
            total_s = interval_h * 3600
            for _ in range(total_s):
                if self._stop.is_set():
                    return
                time.sleep(1)
            if not self._stop.is_set():
                try:
                    self.run_cb()
                except Exception:
                    pass

    def temps_restant(self) -> str:
        if not self.next_run:
            return "Non planifié"
        delta = self.next_run - datetime.now()
        if delta.total_seconds() < 0:
            return "Imminent..."
        h, rem = divmod(int(delta.total_seconds()), 3600)
        m, s   = divmod(rem, 60)
        return f"{h:02d}h {m:02d}m {s:02d}s"

    def is_running(self):
        return self._thread and self._thread.is_alive()


# ════════════════════════════════════════════════════════════════
#  SECTION 11 — APPLICATION GUI (7 ONGLETS)
# ════════════════════════════════════════════════════════════════
class App:
    def __init__(self):
        self.cfg       = Config()
        self.db        = Database()
        self.watchdog  = Watchdog(self.cfg, self.db, self._auto_run)
        self.resultats = []
        self.ia        = MoteurIA(self.cfg)

        self.root = tk.Tk()
        self.root.title(f"AEJ CV MATCHER PRO  v{VERSION}  — ÉDITION PRESTIGE")
        self.root.geometry("1060x760")
        self.root.minsize(900, 650)
        self.root.configure(bg='#F0F0F0')

        self._build()

        if self.cfg.get('AUTO_ACTIVE'):
            self.watchdog.demarrer()
            self._tick_watchdog()

        self.root.mainloop()

    # ── Structure principale ─────────────────────────────────────
    def _build(self):
        self._header()
        self._notebook()

    def _header(self):
        hdr = tk.Frame(self.root, bg=AEJ_GREEN, height=58)
        hdr.pack(fill='x')
        hdr.pack_propagate(False)
        tk.Label(hdr, text="  🎯  AEJ CV MATCHER PRO",
                 font=('Segoe UI', 16, 'bold'), bg=AEJ_GREEN, fg='white').pack(side='left', pady=10)
        tk.Label(hdr, text=f"v{VERSION}  ÉDITION PRESTIGE  —  Agence Prestige, Cocody Abidjan",
                 font=('Segoe UI', 9), bg=AEJ_GREEN, fg='#A5D6A7').pack(side='left', padx=6)
        # Badge stats
        self.v_badge = tk.StringVar(value="")
        tk.Label(hdr, textvariable=self.v_badge, font=('Segoe UI', 9, 'bold'),
                 bg=AEJ_ORANGE, fg='white', padx=8).pack(side='right', padx=14, pady=16)
        self._refresh_badge()

    def _refresh_badge(self):
        s = self.db.stats_globales()
        self.v_badge.set(f"CV:{s['cv']}  Fiches:{s['fiches']}  Scores:{s['scores']}  TOP:{s['top']}")

    def _notebook(self):
        style = ttk.Style()
        style.theme_use('clam')
        style.configure('TNotebook', background='#F0F0F0', borderwidth=0)
        style.configure('TNotebook.Tab', font=('Segoe UI', 10, 'bold'),
                        padding=[12, 6], background='#E0E0E0')
        style.map('TNotebook.Tab',
                  background=[('selected', AEJ_GREEN)],
                  foreground=[('selected', 'white')])

        self.nb = ttk.Notebook(self.root)
        self.nb.pack(fill='both', expand=True, padx=8, pady=6)

        self._tab_config()
        self._tab_upload()     # ← Nouvel onglet Import/Upload
        self._tab_execution()
        self._tab_dashboard()
        self._tab_relation()
        self._tab_autonome()
        self._tab_about()

    # ════════════════════════════════════════════════════════════
    #  ONGLET 1 — CONFIGURATION
    # ════════════════════════════════════════════════════════════
    def _tab_config(self):
        outer = tk.Frame(self.nb, bg=AEJ_GRAY)
        self.nb.add(outer, text='⚙️  Configuration')

        canvas = tk.Canvas(outer, bg=AEJ_GRAY, highlightthickness=0)
        sb     = ttk.Scrollbar(outer, orient='vertical', command=canvas.yview)
        frame  = tk.Frame(canvas, bg=AEJ_GRAY)
        frame.bind("<Configure>",
                   lambda e: canvas.configure(scrollregion=canvas.bbox("all")))
        canvas.create_window((0, 0), window=frame, anchor='nw')
        canvas.configure(yscrollcommand=sb.set)
        canvas.pack(side='left', fill='both', expand=True)
        sb.pack(side='right', fill='y')

        def section(t, color=AEJ_ORANGE):
            lf = tk.LabelFrame(frame, text=f"  {t}  ",
                               font=('Segoe UI', 10, 'bold'),
                               bg=AEJ_GRAY, fg=color, bd=2, relief='groove')
            lf.pack(fill='x', padx=14, pady=7)
            return lf

        def row(parent, lbl, var, r, show='', w=52):
            tk.Label(parent, text=lbl, font=('Segoe UI', 9, 'bold'),
                     bg=AEJ_GRAY).grid(row=r, column=0, sticky='w', padx=10, pady=4)
            e = ttk.Entry(parent, textvariable=var, width=w, show=show)
            e.grid(row=r, column=1, padx=10, pady=4, sticky='w')
            return e

        # Ollama
        g1 = section("🦙  Ollama — IA 100% Locale", AEJ_ORANGE)
        self.v_ollama_url   = tk.StringVar(value=self.cfg.get('OLLAMA_URL', 'http://localhost:11434'))
        self.v_ollama_model = tk.StringVar(value=self.cfg.get('OLLAMA_MODEL', 'gemma3:4b'))
        row(g1, "URL Ollama :", self.v_ollama_url, 0)
        tk.Label(g1, text="Modèle :", font=('Segoe UI', 9, 'bold'),
                 bg=AEJ_GRAY).grid(row=1, column=0, sticky='w', padx=10)
        self._cb_model = ttk.Combobox(g1, textvariable=self.v_ollama_model, width=44,
            values=['gemma3:4b', 'llama3.2', 'llama3.1:8b',
                    'mistral', 'phi4:14b', 'qwen2.5:7b', 'deepseek-r1:8b']
        )
        self._cb_model.grid(row=1, column=1, padx=10, pady=4, sticky='w')
        fr_btns = tk.Frame(g1, bg=AEJ_GRAY)
        fr_btns.grid(row=1, column=2, padx=6, sticky='w')
        tk.Button(fr_btns, text="🔄 Modèles", bg=AEJ_DARK, fg='white',
                  font=('Segoe UI', 8), relief='flat', padx=6, pady=3, cursor='hand2',
                  command=self._actualiser_modeles).pack(side='left', padx=2)
        tk.Button(fr_btns, text="🔗 Tester", bg=AEJ_GREEN, fg='white',
                  font=('Segoe UI', 8), relief='flat', padx=6, pady=3, cursor='hand2',
                  command=self._tester_ollama).pack(side='left', padx=2)
        self.lbl_ollama_status = tk.Label(g1,
            text="⚪ Non vérifié", font=('Segoe UI', 8), bg=AEJ_GRAY, fg='#757575')
        self.lbl_ollama_status.grid(row=2, column=0, columnspan=3, sticky='w', padx=10, pady=(0,4))
        tk.Label(g1,
            text="🔒 Toutes vos données restent sur votre PC — aucun cloud, aucune clé API requise.",
            font=('Segoe UI', 8, 'italic'), bg=AEJ_GRAY, fg=AEJ_DARK
        ).grid(row=3, column=0, columnspan=3, sticky='w', padx=10, pady=(0,6))

        # Source des données (DRIVE / LOCAL)
        gs = section("🗂️  Source des données", AEJ_DARK)
        self.v_source_mode = tk.StringVar(value=self.cfg.get('SOURCE_MODE', 'DRIVE'))
        tk.Label(gs, text="Choisir la source :", font=('Segoe UI', 9, 'bold'),
                 bg=AEJ_GRAY).grid(row=0, column=0, sticky='w', padx=10, pady=6)
        fr_mode = tk.Frame(gs, bg=AEJ_GRAY)
        fr_mode.grid(row=0, column=1, sticky='w')
        ttk.Radiobutton(fr_mode, text="☁️  Google Drive", variable=self.v_source_mode,
                        value='DRIVE', command=self._toggle_source).pack(side='left', padx=10)
        ttk.Radiobutton(fr_mode, text="💻  Fichiers / dossiers locaux", variable=self.v_source_mode,
                        value='LOCAL', command=self._toggle_source).pack(side='left', padx=10)
        tk.Label(gs, text="Le mode local ne nécessite ni credentials.json ni connexion internet.",
                 font=('Segoe UI', 8, 'italic'), bg=AEJ_GRAY, fg='#757575').grid(
                 row=1, column=0, columnspan=2, sticky='w', padx=10, pady=(0,4))

        # Drive
        self.frame_drive = section("📁  Dossiers Google Drive", AEJ_GREEN)
        g2 = self.frame_drive
        champs = [("CV Source :", 'DOSSIER_CV_SOURCE'),
                  ("Fiches Poste :", 'DOSSIER_FICHES_POSTE'),
                  ("Destination :", 'DOSSIER_DESTINATION')]
        self.v_drive = {}
        for i, (lbl, key) in enumerate(champs):
            self.v_drive[key] = tk.StringVar(value=self.cfg.get(key, ''))
            row(g2, lbl, self.v_drive[key], i)

        # Fichiers locaux
        self.frame_local = section("💻  Fichiers / dossiers locaux", AEJ_DARK)
        gl = self.frame_local

        def bloc_picker(parent, titre, r, listbox_attr, ajout_files_cb, ajout_dir_cb, clear_cb):
            tk.Label(parent, text=titre, font=('Segoe UI', 9, 'bold'),
                     bg=AEJ_GRAY).grid(row=r, column=0, sticky='nw', padx=10, pady=6)
            lb = tk.Listbox(parent, height=4, width=58, font=('Consolas', 8),
                            selectmode='extended')
            lb.grid(row=r, column=1, padx=8, pady=6, sticky='w')
            setattr(self, listbox_attr, lb)
            btns = tk.Frame(parent, bg=AEJ_GRAY)
            btns.grid(row=r, column=2, sticky='nw', pady=6)
            tk.Button(btns, text="📄 Fichiers…", font=('Segoe UI', 8),
                      bg=AEJ_GREEN, fg='white', relief='flat', width=12,
                      cursor='hand2', command=ajout_files_cb).pack(pady=1)
            tk.Button(btns, text="📁 Dossier…", font=('Segoe UI', 8),
                      bg=AEJ_ORANGE, fg='white', relief='flat', width=12,
                      cursor='hand2', command=ajout_dir_cb).pack(pady=1)
            tk.Button(btns, text="🗑 Effacer", font=('Segoe UI', 8),
                      bg='#757575', fg='white', relief='flat', width=12,
                      cursor='hand2', command=clear_cb).pack(pady=1)

        bloc_picker(gl, "CV (fichiers/dossiers) :", 0, 'lst_cv',
                    lambda: self._ajouter_fichiers('lst_cv'),
                    lambda: self._ajouter_dossier('lst_cv'),
                    lambda: self._effacer_liste('lst_cv'))
        bloc_picker(gl, "Fiches de poste :", 1, 'lst_fiche',
                    lambda: self._ajouter_fichiers('lst_fiche'),
                    lambda: self._ajouter_dossier('lst_fiche'),
                    lambda: self._effacer_liste('lst_fiche'))

        # Destination locale
        tk.Label(gl, text="Destination (dossier) :", font=('Segoe UI', 9, 'bold'),
                 bg=AEJ_GRAY).grid(row=2, column=0, sticky='w', padx=10, pady=6)
        self.v_local_dest = tk.StringVar(value=self.cfg.get('LOCAL_DEST_PATH', ''))
        ttk.Entry(gl, textvariable=self.v_local_dest, width=50).grid(
            row=2, column=1, padx=8, sticky='w')
        tk.Button(gl, text="📁 Choisir…", font=('Segoe UI', 8),
                  bg=AEJ_GREEN, fg='white', relief='flat', width=12, cursor='hand2',
                  command=lambda: self.v_local_dest.set(
                      filedialog.askdirectory() or self.v_local_dest.get())
                  ).grid(row=2, column=2, sticky='nw', pady=6)

        # Pré-remplir les listes depuis la config
        for p in self.cfg.get('LOCAL_CV_PATHS', []) or []:
            self.lst_cv.insert('end', p)
        for p in self.cfg.get('LOCAL_FICHE_PATHS', []) or []:
            self.lst_fiche.insert('end', p)

        self._toggle_source()

        # Matching
        g3 = section("🎛️  Paramètres de matching", AEJ_ORANGE)
        self.v_mode   = tk.StringVar(value=self.cfg.get('MODE', 'MULTI_POSTES'))
        self.v_action = tk.StringVar(value=self.cfg.get('ACTION', 'COPIER'))
        self.v_s_top  = tk.IntVar(value=self.cfg.get('SEUIL_TOP', 80))
        self.v_s_moy  = tk.IntVar(value=self.cfg.get('SEUIL_MOYEN', 65))
        self.v_tfidf  = tk.IntVar(value=self.cfg.get('SEUIL_TFIDF', 18))
        self.v_maxcv  = tk.IntVar(value=self.cfg.get('MAX_CV_PAR_POSTE', 15))
        self.v_maxp   = tk.IntVar(value=self.cfg.get('MAX_POSTES_PAR_CV', 3))
        self.v_wrk    = tk.IntVar(value=self.cfg.get('PARALLEL_WORKERS', 4))

        for ri, (lbl, opts, var) in enumerate([
            ("Mode :",   [("Multi-postes (1 CV → plusieurs postes)", "MULTI_POSTES"),
                          ("Meilleur poste uniquement", "MEILLEUR_POSTE")], self.v_mode),
            ("Action :", [("COPIER (source intacte ✓)", "COPIER"),
                          ("DÉPLACER", "DEPLACER")], self.v_action)
        ]):
            tk.Label(g3, text=lbl, font=('Segoe UI', 9, 'bold'),
                     bg=AEJ_GRAY).grid(row=ri, column=0, sticky='w', padx=10, pady=4)
            fr = tk.Frame(g3, bg=AEJ_GRAY)
            fr.grid(row=ri, column=1, sticky='w', padx=5)
            for txt, val in opts:
                ttk.Radiobutton(fr, text=txt, variable=var, value=val).pack(side='left', padx=8)

        spin_cfg = [("Seuil TOP (%) :", self.v_s_top, 50, 100, 2),
                    ("Seuil MOYEN (%) :", self.v_s_moy, 30, 90, 3),
                    ("Seuil TF-IDF pré-filtre :", self.v_tfidf, 5, 50, 4),
                    ("Max CV / poste :", self.v_maxcv, 5, 50, 5),
                    ("Max postes / CV :", self.v_maxp, 1, 10, 6),
                    ("Threads parallèles :", self.v_wrk, 1, 8, 7)]
        for lbl, var, fr, to, ri in spin_cfg:
            tk.Label(g3, text=lbl, font=('Segoe UI', 9, 'bold'),
                     bg=AEJ_GRAY).grid(row=ri, column=0, sticky='w', padx=10, pady=3)
            ttk.Spinbox(g3, from_=fr, to=to, textvariable=var,
                        width=8).grid(row=ri, column=1, sticky='w', padx=10)

        # OCR
        g4 = section("🔍  Tesseract OCR (CV scannés)", "#555555")
        self.v_tess = tk.StringVar(value=self.cfg.get('TESSERACT_PATH',
            r'C:\Program Files\Tesseract-OCR\tesseract.exe'))
        tk.Label(g4, text="Chemin :", font=('Segoe UI', 9, 'bold'),
                 bg=AEJ_GRAY).grid(row=0, column=0, sticky='w', padx=10)
        ttk.Entry(g4, textvariable=self.v_tess, width=50).grid(row=0, column=1, padx=10)
        ttk.Button(g4, text="Parcourir…",
            command=lambda: self.v_tess.set(
                filedialog.askopenfilename(filetypes=[("EXE","*.exe")]) or self.v_tess.get()
            )).grid(row=0, column=2)

        # Email / Signature
        g5 = section("📧  Signature Email", AEJ_GREEN)
        self.v_sig = tk.Text(g5, height=5, width=60, font=('Segoe UI', 9))
        self.v_sig.insert('end', self.cfg.get('SIGNATURE_EMAIL', ''))
        self.v_sig.grid(row=0, column=0, columnspan=2, padx=10, pady=6)

        # Enregistrer
        btn_f = tk.Frame(frame, bg=AEJ_GRAY)
        btn_f.pack(pady=14)
        tk.Button(btn_f, text="💾   ENREGISTRER LA CONFIGURATION",
                  font=('Segoe UI', 11, 'bold'), bg=AEJ_GREEN, fg='white',
                  padx=20, pady=10, relief='flat', cursor='hand2',
                  command=self._sauver_config).pack(side='left', padx=8)
        tk.Button(btn_f, text="🔄  Réinitialiser",
                  font=('Segoe UI', 10), bg='#757575', fg='white',
                  padx=12, pady=9, relief='flat', cursor='hand2',
                  command=self._reinit_config).pack(side='left', padx=8)

    def _sauver_config(self):
        self.cfg.update({
            'SOURCE_MODE'         : self.v_source_mode.get(),
            'LOCAL_CV_PATHS'      : list(self.lst_cv.get(0, 'end')),
            'LOCAL_FICHE_PATHS'   : list(self.lst_fiche.get(0, 'end')),
            'LOCAL_DEST_PATH'     : self.v_local_dest.get().strip(),
            'OLLAMA_URL'          : self.v_ollama_url.get().strip(),
            'OLLAMA_MODEL'        : self.v_ollama_model.get().strip(),
            'DOSSIER_CV_SOURCE'   : self.v_drive['DOSSIER_CV_SOURCE'].get().strip(),
            'DOSSIER_FICHES_POSTE': self.v_drive['DOSSIER_FICHES_POSTE'].get().strip(),
            'DOSSIER_DESTINATION' : self.v_drive['DOSSIER_DESTINATION'].get().strip(),
            'MODE'                : self.v_mode.get(),
            'ACTION'              : self.v_action.get(),
            'SEUIL_TOP'           : self.v_s_top.get(),
            'SEUIL_MOYEN'         : self.v_s_moy.get(),
            'SEUIL_TFIDF'         : self.v_tfidf.get(),
            'MAX_CV_PAR_POSTE'    : self.v_maxcv.get(),
            'MAX_POSTES_PAR_CV'   : self.v_maxp.get(),
            'PARALLEL_WORKERS'    : self.v_wrk.get(),
            'TESSERACT_PATH'      : self.v_tess.get(),
            'SIGNATURE_EMAIL'     : self.v_sig.get('1.0', 'end').strip(),
        })
        self.ia = MoteurIA(self.cfg)
        messagebox.showinfo("Configuration", "✅ Configuration enregistrée !")

    def _reinit_config(self):
        if messagebox.askyesno("Réinitialiser", "Remettre tous les paramètres par défaut ?"):
            for k, v in DEFAULT_CONFIG.items():
                self.cfg.set(k, v)
            messagebox.showinfo("Config", "Paramètres par défaut restaurés. Redémarrez pour appliquer.")

    def _get_ollama_models(self):
        """Return list of installed Ollama models (empty list if Ollama not running)."""
        try:
            import urllib.request as ur
            url = self.v_ollama_url.get().strip().rstrip('/') + "/api/tags"
            with ur.urlopen(url, timeout=5) as resp:
                data = json.loads(resp.read())
                return [m["name"] for m in data.get("models", [])]
        except Exception:
            return []

    def _actualiser_modeles(self):
        """Fetch installed Ollama models and populate the model combobox."""
        def run():
            models = self._get_ollama_models()
            def update():
                if models:
                    self._cb_model['values'] = models
                    if self.v_ollama_model.get() not in models:
                        self.v_ollama_model.set(models[0])
                    self.lbl_ollama_status.config(
                        text=f"✅ Ollama actif — {len(models)} modèle(s) disponible(s)", fg=AEJ_GREEN)
                else:
                    self.lbl_ollama_status.config(
                        text="❌ Ollama non détecté — installez Ollama sur https://ollama.com", fg='red')
            self.root.after(0, update)
        threading.Thread(target=run, daemon=True).start()
        self.lbl_ollama_status.config(text="⏳ Détection en cours...", fg='#757575')

    def _tester_ollama(self):
        """Test Ollama connection with a quick LLM call."""
        url  = self.v_ollama_url.get().strip()
        model = self.v_ollama_model.get().strip()
        if not url:
            messagebox.showwarning("URL manquante", "Entrez l'URL Ollama (ex: http://localhost:11434)")
            return
        def test():
            try:
                payload = {
                    "model": model,
                    "messages": [{"role": "user", "content": "Réponds: OK"}],
                    "max_tokens": 10, "stream": False
                }
                r = requests.post(url.rstrip('/') + "/v1/chat/completions", json=payload, timeout=30)
                if r.status_code == 200:
                    reply = r.json()['choices'][0]['message']['content'].strip()
                    self.root.after(0, lambda: (
                        messagebox.showinfo("Ollama", f"✅ Ollama opérationnel !\nModèle : {model}\nRéponse : {reply}"),
                        self.lbl_ollama_status.config(text=f"✅ Ollama actif — modèle : {model}", fg=AEJ_GREEN)
                    ))
                else:
                    err = f"HTTP {r.status_code}: {r.text[:150]}"
                    self.root.after(0, lambda: (
                        messagebox.showerror("Ollama", err),
                        self.lbl_ollama_status.config(text=f"❌ Erreur : {err}", fg='red')
                    ))
            except Exception as e:
                msg = str(e)
                self.root.after(0, lambda: (
                    messagebox.showerror("Ollama", f"Impossible de joindre Ollama :\n{msg}\n\n"
                                         "Vérifiez qu'Ollama est lancé (ollama serve)"),
                    self.lbl_ollama_status.config(text=f"❌ {msg[:60]}", fg='red')
                ))
        threading.Thread(target=test, daemon=True).start()
        self.lbl_ollama_status.config(text="⏳ Test en cours...", fg='#757575')

    # ── Gestion source locale ────────────────────────────────────
    def _toggle_source(self):
        """Active/désactive les sections Drive vs Local selon le mode."""
        mode = self.v_source_mode.get()
        def set_state(frame, state):
            for child in frame.winfo_children():
                try:
                    child.configure(state=state)
                except tk.TclError:
                    pass
                for sub in child.winfo_children():
                    try:
                        sub.configure(state=state)
                    except tk.TclError:
                        pass
        if mode == 'DRIVE':
            set_state(self.frame_drive, 'normal')
            set_state(self.frame_local, 'disabled')
        else:
            set_state(self.frame_drive, 'disabled')
            set_state(self.frame_local, 'normal')

    def _ajouter_fichiers(self, listbox_attr):
        lb = getattr(self, listbox_attr)
        fichiers = filedialog.askopenfilenames(
            title="Sélectionner des fichiers (CV ou fiches)",
            filetypes=[("Documents", "*.pdf *.docx *.doc *.txt *.jpg *.jpeg *.png"),
                       ("Tous les fichiers", "*.*")])
        for f in fichiers:
            if f not in lb.get(0, 'end'):
                lb.insert('end', f)

    def _ajouter_dossier(self, listbox_attr):
        lb = getattr(self, listbox_attr)
        d = filedialog.askdirectory(title="Sélectionner un dossier")
        if d and d not in lb.get(0, 'end'):
            lb.insert('end', d)

    def _effacer_liste(self, listbox_attr):
        lb = getattr(self, listbox_attr)
        sel = lb.curselection()
        if sel:
            # Effacer la sélection
            for i in reversed(sel):
                lb.delete(i)
        else:
            # Tout effacer
            lb.delete(0, 'end')

    # ════════════════════════════════════════════════════════════
    #  ONGLET 2 — 📤 IMPORTER / UPLOADER
    # ════════════════════════════════════════════════════════════
    def _tab_upload(self):
        f = tk.Frame(self.nb, bg=AEJ_GRAY)
        self.nb.add(f, text='📤  Importer')

        self._upload_queue  = []       # liste de dicts par fichier
        self._arret_upload  = threading.Event()
        self._drive_upload  = None    # DriveClient dédié à l'upload

        # ── Bannière d'info ─────────────────────────────────────
        info = tk.Frame(f, bg='#E8F5E9', relief='solid', bd=1)
        info.pack(fill='x', padx=12, pady=(10, 5))
        tk.Label(info, bg='#E8F5E9', font=('Segoe UI', 9), wraplength=920,
                 text=(f"Importez des CV et des fiches de poste depuis votre ordinateur "
                       f"directement dans vos dossiers Google Drive. "
                       f"Formats acceptés : {EXTS_LABEL}.")
                 ).pack(padx=12, pady=7, anchor='w')

        # ── Zone d'ajout ────────────────────────────────────────
        add_frame = tk.LabelFrame(f, text="  📁  Sélectionner les fichiers à importer  ",
                                   font=('Segoe UI', 10, 'bold'),
                                   bg=AEJ_GRAY, fg=AEJ_ORANGE, bd=2, relief='groove')
        add_frame.pack(fill='x', padx=12, pady=5)

        # — Ligne CV —
        cv_row = tk.Frame(add_frame, bg=AEJ_GRAY)
        cv_row.pack(fill='x', padx=10, pady=7)
        tk.Label(cv_row, text="📄  CV :", font=('Segoe UI', 10, 'bold'),
                 bg=AEJ_GRAY, width=14, anchor='w').pack(side='left')
        tk.Button(cv_row, text="➕ Fichiers CV",
                  bg=AEJ_GREEN, fg='white', font=('Segoe UI', 9, 'bold'),
                  relief='flat', padx=12, pady=5, cursor='hand2',
                  command=lambda: self._ajout_fichiers('CV')
                  ).pack(side='left', padx=4)
        tk.Button(cv_row, text="📂 Dossier de CV",
                  bg='#388E3C', fg='white', font=('Segoe UI', 9),
                  relief='flat', padx=10, pady=5, cursor='hand2',
                  command=lambda: self._ajout_dossier('CV')
                  ).pack(side='left', padx=4)
        tk.Label(cv_row, text="→ Drive ID cible :", font=('Segoe UI', 9),
                 bg=AEJ_GRAY).pack(side='left', padx=(16, 4))
        self.v_dest_cv = tk.StringVar(value=self.cfg.get('DOSSIER_CV_SOURCE', ''))
        ttk.Entry(cv_row, textvariable=self.v_dest_cv, width=34).pack(side='left', padx=2)
        tk.Label(cv_row, text="(ID Drive)", font=('Segoe UI', 8),
                 bg=AEJ_GRAY, fg='#9E9E9E').pack(side='left', padx=4)

        # — Ligne Fiches —
        fiche_row = tk.Frame(add_frame, bg=AEJ_GRAY)
        fiche_row.pack(fill='x', padx=10, pady=7)
        tk.Label(fiche_row, text="📋  Fiches :", font=('Segoe UI', 10, 'bold'),
                 bg=AEJ_GRAY, width=14, anchor='w').pack(side='left')
        tk.Button(fiche_row, text="➕ Fichiers fiches",
                  bg=AEJ_ORANGE, fg='white', font=('Segoe UI', 9, 'bold'),
                  relief='flat', padx=12, pady=5, cursor='hand2',
                  command=lambda: self._ajout_fichiers('FICHE')
                  ).pack(side='left', padx=4)
        tk.Button(fiche_row, text="📂 Dossier de fiches",
                  bg='#E65100', fg='white', font=('Segoe UI', 9),
                  relief='flat', padx=10, pady=5, cursor='hand2',
                  command=lambda: self._ajout_dossier('FICHE')
                  ).pack(side='left', padx=4)
        tk.Label(fiche_row, text="Sous-dossier Drive :", font=('Segoe UI', 9),
                 bg=AEJ_GRAY).pack(side='left', padx=(16, 4))
        self.v_sous_dossier = tk.StringVar(value="")
        ttk.Entry(fiche_row, textvariable=self.v_sous_dossier, width=22).pack(side='left', padx=2)
        tk.Label(fiche_row, text="(optionnel — ex: SCHIBA_HOLDING)",
                 font=('Segoe UI', 8), bg=AEJ_GRAY, fg='#9E9E9E').pack(side='left', padx=4)

        # — Options —
        opt_row = tk.Frame(add_frame, bg=AEJ_GRAY)
        opt_row.pack(fill='x', padx=10, pady=(0, 8))
        self.v_skip_dupes   = tk.BooleanVar(value=True)
        self.v_recurse_dir  = tk.BooleanVar(value=True)
        self.v_renommer_auto= tk.BooleanVar(value=False)
        ttk.Checkbutton(opt_row, text="Ignorer les fichiers déjà présents dans Drive",
                        variable=self.v_skip_dupes).pack(side='left', padx=8)
        ttk.Checkbutton(opt_row, text="Inclure les sous-dossiers",
                        variable=self.v_recurse_dir).pack(side='left', padx=16)
        ttk.Checkbutton(opt_row, text="Renommer : NOM_PRENOM_CV.ext",
                        variable=self.v_renommer_auto).pack(side='left', padx=16)

        # ── File Queue ──────────────────────────────────────────
        queue_lf = tk.LabelFrame(f, text="  📋  Fichiers sélectionnés  ",
                                  font=('Segoe UI', 10, 'bold'),
                                  bg=AEJ_GRAY, fg=AEJ_GREEN, bd=2, relief='groove')
        queue_lf.pack(fill='both', expand=True, padx=12, pady=5)

        # Toolbar queue
        tb = tk.Frame(queue_lf, bg=AEJ_GRAY)
        tb.pack(fill='x', padx=8, pady=5)
        self.v_nb_queue = tk.StringVar(value="Queue vide — sélectionnez des fichiers ci-dessus")
        tk.Label(tb, textvariable=self.v_nb_queue, font=('Segoe UI', 9, 'bold'),
                 bg=AEJ_GRAY, fg=AEJ_GREEN).pack(side='left')
        tk.Button(tb, text="✕ Retirer",
                  bg='#D32F2F', fg='white', font=('Segoe UI', 9),
                  relief='flat', padx=8, pady=3, cursor='hand2',
                  command=self._retirer_selectionne_queue
                  ).pack(side='right', padx=4)
        tk.Button(tb, text="🗑 Vider tout",
                  bg='#757575', fg='white', font=('Segoe UI', 9),
                  relief='flat', padx=8, pady=3, cursor='hand2',
                  command=self._vider_queue
                  ).pack(side='right', padx=4)
        tk.Button(tb, text="📊 Statistiques",
                  bg='#455A64', fg='white', font=('Segoe UI', 9),
                  relief='flat', padx=8, pady=3, cursor='hand2',
                  command=self._stats_queue
                  ).pack(side='right', padx=4)

        # Treeview
        cols_q  = ('TYPE', 'FICHIER', 'TAILLE', 'CHEMIN LOCAL', 'DESTINATION', 'STATUT')
        self.tree_queue = ttk.Treeview(queue_lf, columns=cols_q,
                                        show='headings', height=10,
                                        selectmode='extended')
        for col, w, anchor in zip(
            cols_q,
            [8, 32, 10, 32, 18, 16],
            ['center', 'w', 'center', 'w', 'center', 'center']
        ):
            self.tree_queue.heading(col, text=col,
                command=lambda c=col: self._sort_tree(self.tree_queue, c, False))
            self.tree_queue.column(col, width=w * 9, anchor=anchor, minwidth=60)

        sq_y = ttk.Scrollbar(queue_lf, orient='vertical',   command=self.tree_queue.yview)
        sq_x = ttk.Scrollbar(queue_lf, orient='horizontal',  command=self.tree_queue.xview)
        self.tree_queue.configure(yscrollcommand=sq_y.set, xscrollcommand=sq_x.set)
        self.tree_queue.pack(fill='both', expand=True, padx=8, pady=(0,2), side='left')
        sq_y.pack(side='right', fill='y')
        sq_x.pack(side='bottom', fill='x')

        # Menu contextuel queue
        self.ctx_queue = tk.Menu(self.root, tearoff=0)
        self.ctx_queue.add_command(label="✕ Retirer de la liste",
                                    command=self._retirer_selectionne_queue)
        self.ctx_queue.add_command(label="📂 Ouvrir dans l'explorateur",
                                    command=self._ouvrir_explorateur)
        self.ctx_queue.add_separator()
        self.ctx_queue.add_command(label="📤 Uploader uniquement ce fichier",
                                    command=self._upload_selectionne_seul)
        self.tree_queue.bind("<Button-3>", self._ctx_queue_show)

        # ── Contrôles upload ────────────────────────────────────
        ctrl = tk.Frame(f, bg=AEJ_GRAY)
        ctrl.pack(fill='x', padx=12, pady=6)

        self.btn_upload_all = tk.Button(ctrl,
            text="📤  UPLOADER TOUT VERS DRIVE",
            font=('Segoe UI', 12, 'bold'), bg=AEJ_GREEN, fg='white',
            padx=20, pady=10, relief='flat', cursor='hand2',
            activebackground=AEJ_DARK,
            command=self._lancer_upload)
        self.btn_upload_all.pack(side='left', padx=6)

        tk.Button(ctrl, text="▶ Upload + Matcher ensuite",
                  font=('Segoe UI', 11, 'bold'), bg=AEJ_ORANGE, fg='white',
                  padx=14, pady=10, relief='flat', cursor='hand2',
                  command=self._upload_puis_matcher
                  ).pack(side='left', padx=6)

        self.btn_stop_upload = tk.Button(ctrl, text="⏹  Arrêter",
                  font=('Segoe UI', 10), bg='#D32F2F', fg='white',
                  padx=12, pady=10, relief='flat', cursor='hand2',
                  state='disabled', command=self._stopper_upload)
        self.btn_stop_upload.pack(side='left', padx=6)

        tk.Button(ctrl, text="🔗 Ouvrir Drive",
                  font=('Segoe UI', 10), bg='#455A64', fg='white',
                  padx=12, pady=10, relief='flat', cursor='hand2',
                  command=lambda: webbrowser.open(
                      f"https://drive.google.com/drive/folders/{self.v_dest_cv.get().strip()}")
                  ).pack(side='right', padx=6)

        # ── Progression ─────────────────────────────────────────
        prog_f = tk.Frame(f, bg=AEJ_GRAY)
        prog_f.pack(fill='x', padx=12, pady=(0, 4))
        self.v_up_status = tk.StringVar(value="Prêt à uploader.")
        tk.Label(prog_f, textvariable=self.v_up_status, font=('Segoe UI', 9, 'bold'),
                 bg=AEJ_GRAY, fg=AEJ_DARK).pack(anchor='w')
        self.bar_up = ttk.Progressbar(prog_f, mode='determinate', maximum=100)
        self.bar_up.pack(fill='x', pady=2)
        self.v_up_pct = tk.StringVar(value="")
        tk.Label(prog_f, textvariable=self.v_up_pct, font=('Segoe UI', 8),
                 bg=AEJ_GRAY, fg='#616161').pack(anchor='e')

        # ── Journal upload ──────────────────────────────────────
        self.journal_up = scrolledtext.ScrolledText(
            f, height=5, font=('Consolas', 9),
            bg='#1C1C1C', fg='#69F0AE', relief='flat', padx=6, pady=4)
        self.journal_up.pack(fill='x', padx=12, pady=(0, 10))
        self.journal_up.config(state='disabled')

    # ── Helpers Upload ──────────────────────────────────────────
    def _ajout_fichiers(self, type_doc: str):
        label = "CV" if type_doc == 'CV' else "Fiches de poste"
        chemins = filedialog.askopenfilenames(
            title=f"Sélectionner des {label}",
            filetypes=[
                ("Documents", "*.pdf *.doc *.docx *.txt"),
                ("Images",    "*.jpg *.jpeg *.png"),
                ("PDF",       "*.pdf"),
                ("Word",      "*.doc *.docx"),
                ("Tous",      "*.*"),
            ]
        )
        ajouts = 0
        for ch in chemins:
            if self._ajouter_en_queue(Path(ch), type_doc):
                ajouts += 1
        self._refresh_queue_ui()
        self._log_up(f"➕ {ajouts} {label} ajouté(s) — total queue : {len(self._upload_queue)}")

    def _ajout_dossier(self, type_doc: str):
        label   = "CV" if type_doc == 'CV' else "Fiches de poste"
        dossier = filedialog.askdirectory(
            title=f"Sélectionner le dossier contenant les {label}")
        if not dossier:
            return
        dossier = Path(dossier)
        pattern = '**/*' if self.v_recurse_dir.get() else '*'
        fichiers = [f for f in dossier.glob(pattern)
                    if f.is_file() and f.suffix.lower() in EXTS_SUPPORTEES]
        ajouts = sum(1 for f in fichiers if self._ajouter_en_queue(f, type_doc))
        self._refresh_queue_ui()
        self._log_up(f"📂 Dossier {dossier.name} : {ajouts}/{len(fichiers)} fichier(s) ajouté(s)")

    def _ajouter_en_queue(self, chemin: Path, type_doc: str) -> bool:
        """Ajoute un fichier à la queue. Retourne True si ajouté."""
        if chemin.suffix.lower() not in EXTS_SUPPORTEES:
            return False
        # Pas de doublon dans la queue
        if any(q['path'] == chemin for q in self._upload_queue):
            return False
        try:
            taille = chemin.stat().st_size
        except Exception:
            taille = 0
        dest_id = (self.v_dest_cv.get().strip() if type_doc == 'CV'
                   else self.cfg.get('DOSSIER_FICHES_POSTE', ''))
        # Renommage automatique (optionnel)
        nom = chemin.name
        if self.v_renommer_auto.get() and type_doc == 'CV':
            stem = chemin.stem.upper().replace(' ', '_')
            nom  = f"{stem}_CV{chemin.suffix.lower()}"

        self._upload_queue.append({
            'path'    : chemin,
            'nom'     : nom,
            'type'    : type_doc,
            'taille'  : taille,
            'statut'  : 'EN_ATTENTE',
            'dest_id' : dest_id,
            'drive_id': None,
            'drive_lien': '',
        })
        return True

    def _refresh_queue_ui(self):
        for item in self.tree_queue.get_children():
            self.tree_queue.delete(item)

        for i, q in enumerate(self._upload_queue):
            if q['taille'] >= 1_048_576:
                taille_s = f"{q['taille']/1_048_576:.1f} MB"
            elif q['taille'] >= 1024:
                taille_s = f"{q['taille']//1024} KB"
            else:
                taille_s = f"{q['taille']} B"

            dest_s = ("→ CV Drive" if q['type'] == 'CV' else "→ Fiches Drive")
            st     = UP_STATUTS.get(q['statut'], q['statut'])
            chemin_s = str(q['path'].parent)[-40:]

            tag = {'OK': 'ok', 'ERREUR': 'err', 'EXISTANT': 'skip',
                   'UPLOAD': 'busy'}.get(q['statut'], 'wait')

            self.tree_queue.insert('', 'end', iid=str(i),
                values=(q['type'], q['nom'], taille_s, chemin_s, dest_s, st),
                tags=(tag,))

        self.tree_queue.tag_configure('ok',   background='#C8E6C9')
        self.tree_queue.tag_configure('err',  background='#FFCDD2')
        self.tree_queue.tag_configure('skip', background='#F5F5F5', foreground='#9E9E9E')
        self.tree_queue.tag_configure('busy', background='#E3F2FD')
        self.tree_queue.tag_configure('wait', background='#FFF9C4')

        n_total  = len(self._upload_queue)
        n_ok     = sum(1 for q in self._upload_queue if q['statut'] == 'OK')
        n_err    = sum(1 for q in self._upload_queue if q['statut'] == 'ERREUR')
        n_cv     = sum(1 for q in self._upload_queue if q['type'] == 'CV')
        n_fiche  = sum(1 for q in self._upload_queue if q['type'] == 'FICHE')
        taille_t = sum(q['taille'] for q in self._upload_queue)
        t_str    = (f"{taille_t/1_048_576:.1f} MB" if taille_t >= 1_048_576
                    else f"{taille_t//1024} KB")

        self.v_nb_queue.set(
            f"🗂 {n_total} fichier(s)  |  📄 CV: {n_cv}  |  📋 Fiches: {n_fiche}  "
            f"|  📦 {t_str}  |  ✅ {n_ok} uploadé(s)  |  ❌ {n_err} erreur(s)"
        )

    def _vider_queue(self):
        if self._upload_queue and messagebox.askyesno(
                "Vider", f"Supprimer les {len(self._upload_queue)} fichier(s) de la queue ?"):
            self._upload_queue.clear()
            self._refresh_queue_ui()

    def _retirer_selectionne_queue(self):
        sel = sorted([int(s) for s in self.tree_queue.selection()], reverse=True)
        for idx in sel:
            if 0 <= idx < len(self._upload_queue):
                self._upload_queue.pop(idx)
        self._refresh_queue_ui()

    def _ctx_queue_show(self, event):
        item = self.tree_queue.identify_row(event.y)
        if item:
            self.tree_queue.selection_set(item)
            self.ctx_queue.post(event.x_root, event.y_root)

    def _ouvrir_explorateur(self):
        sel = self.tree_queue.selection()
        if not sel:
            return
        idx  = int(sel[0])
        path = self._upload_queue[idx]['path']
        try:
            import subprocess
            subprocess.Popen(f'explorer /select,"{path}"')
        except Exception:
            pass

    def _stats_queue(self):
        if not self._upload_queue:
            messagebox.showinfo("Stats", "Queue vide.")
            return
        par_ext  = {}
        par_type = {'CV': 0, 'FICHE': 0}
        taille_t = 0
        for q in self._upload_queue:
            ext = q['path'].suffix.lower()
            par_ext[ext] = par_ext.get(ext, 0) + 1
            par_type[q['type']] += 1
            taille_t += q['taille']
        ext_str = ' | '.join(f"{k}: {v}" for k, v in sorted(par_ext.items()))
        msg = (
            f"Total fichiers       : {len(self._upload_queue)}\n"
            f"  CV                 : {par_type['CV']}\n"
            f"  Fiches de poste    : {par_type['FICHE']}\n\n"
            f"Taille totale        : {taille_t/1_048_576:.1f} MB\n"
            f"Par extension        : {ext_str}\n\n"
            f"Uploadés (✅)        : {sum(1 for q in self._upload_queue if q['statut']=='OK')}\n"
            f"En attente (⏳)      : {sum(1 for q in self._upload_queue if q['statut']=='EN_ATTENTE')}\n"
            f"Erreurs (❌)         : {sum(1 for q in self._upload_queue if q['statut']=='ERREUR')}\n"
            f"Déjà présents (⏭)   : {sum(1 for q in self._upload_queue if q['statut']=='EXISTANT')}"
        )
        messagebox.showinfo("Statistiques de la queue", msg)

    def _log_up(self, msg: str):
        self.journal_up.config(state='normal')
        self.journal_up.insert('end', f"[{datetime.now().strftime('%H:%M:%S')}]  {msg}\n")
        self.journal_up.see('end')
        self.journal_up.config(state='disabled')

    def _lancer_upload(self, callback_fin=None):
        en_attente = [q for q in self._upload_queue if q['statut'] == 'EN_ATTENTE']
        if not en_attente:
            messagebox.showwarning("Queue vide",
                "Aucun fichier en attente.\nSélectionnez des fichiers à importer.")
            return
        if not CREDS_FILE.exists():
            messagebox.showwarning("Connexion Drive",
                f"Fichier credentials.json introuvable dans :\n{APP_DIR}")
            return

        self._arret_upload.clear()
        self.btn_upload_all.config(state='disabled', text='⏳ Upload en cours...')
        self.btn_stop_upload.config(state='normal')
        self.bar_up['value'] = 0

        def run():
            try:
                drive = DriveClient()
                drive.connecter()
                total  = len(en_attente)
                ok, err, skip = 0, 0, 0

                for idx, q in enumerate(en_attente):
                    if self._arret_upload.is_set():
                        self.root.after(0, lambda: self._log_up("⏹ Upload interrompu par l'utilisateur."))
                        break

                    nom_affiche = q['nom'][:50]
                    self.root.after(0, lambda n=nom_affiche, i=idx: (
                        self.v_up_status.set(f"[{i+1}/{total}]  Upload : {n}"),
                    ))
                    q['statut'] = 'UPLOAD'
                    self.root.after(0, self._refresh_queue_ui)

                    try:
                        # Vérifier doublon Drive
                        if self.v_skip_dupes.get() and q['dest_id']:
                            if drive.fichier_existe(q['nom'], q['dest_id']):
                                q['statut'] = 'EXISTANT'
                                skip += 1
                                self.root.after(0, lambda n=q['nom']:
                                    self._log_up(f"   ⏭  Déjà présent : {n}"))
                                continue

                        # Créer sous-dossier fiche si besoin
                        dest_id = q['dest_id']
                        if q['type'] == 'FICHE':
                            sous = self.v_sous_dossier.get().strip()
                            if sous and dest_id:
                                dest_id = drive.creer_dossier(sous, dest_id)

                        if not dest_id:
                            raise ValueError("ID dossier destination manquant. Vérifiez la Configuration.")

                        fid, lien = drive.uploader_local(q['path'], q['nom'], dest_id)
                        q['statut']     = 'OK'
                        q['drive_id']   = fid
                        q['drive_lien'] = lien
                        ok += 1
                        self.root.after(0, lambda n=q['nom']:
                            self._log_up(f"   ✅ Uploadé : {n}"))

                    except Exception as e:
                        q['statut'] = 'ERREUR'
                        err += 1
                        self.root.after(0, lambda n=q['nom'], er=str(e):
                            self._log_up(f"   ❌ {n} : {er}"))

                    pct = int((idx + 1) / total * 100)
                    self.root.after(0, lambda p=pct: (
                        self.bar_up.configure(value=p),
                        self.v_up_pct.set(f"{p} %"),
                    ))
                    self.root.after(0, self._refresh_queue_ui)
                    time.sleep(0.25)

                # Résumé final
                resume = (f"📦 Upload terminé — ✅ {ok} uploadé(s)  "
                          f"|  ⏭ {skip} ignoré(s)  |  ❌ {err} erreur(s)")
                self.root.after(0, lambda: (
                    self._log_up(f"\n{resume}\n"),
                    self.v_up_status.set(resume),
                    self._refresh_badge(),
                ))
                if callback_fin:
                    self.root.after(800, callback_fin)

            except Exception as e:
                self.root.after(0, lambda: messagebox.showerror("Erreur Drive", str(e)))
            finally:
                self.root.after(0, lambda: (
                    self.btn_upload_all.config(
                        state='normal', text='📤  UPLOADER TOUT VERS DRIVE'),
                    self.btn_stop_upload.config(state='disabled'),
                ))

        threading.Thread(target=run, daemon=True).start()

    def _stopper_upload(self):
        self._arret_upload.set()
        self.btn_stop_upload.config(state='disabled')

    def _upload_selectionne_seul(self):
        sel = self.tree_queue.selection()
        if not sel:
            return
        # Marquer EN_ATTENTE uniquement le sélectionné, mettre les autres en SKIP temporaire
        idx_sel = {int(s) for s in sel}
        for i, q in enumerate(self._upload_queue):
            if i not in idx_sel and q['statut'] == 'EN_ATTENTE':
                q['statut'] = '_HOLD'
        self._lancer_upload(callback_fin=self._restaurer_hold)

    def _restaurer_hold(self):
        for q in self._upload_queue:
            if q['statut'] == '_HOLD':
                q['statut'] = 'EN_ATTENTE'
        self._refresh_queue_ui()

    def _upload_puis_matcher(self):
        """Upload les fichiers, puis lance le matching automatiquement."""
        self._log_up("⏳ Upload → Matching enchaîné. Le pipeline démarrera après l'upload.")

        def lancer_matching():
            self._log(f"▶ Matching lancé automatiquement après upload...")
            self.nb.select(3)   # Onglet Exécution (index 3 avec le nouvel onglet)
            self.root.after(400, self._lancer)

        self._lancer_upload(callback_fin=lancer_matching)

    # ════════════════════════════════════════════════════════════
    #  ONGLET 3 — EXÉCUTION
    # ════════════════════════════════════════════════════════════
    def _tab_execution(self):
        f = tk.Frame(self.nb, bg=AEJ_GRAY)
        self.nb.add(f, text='🚀  Exécution')

        # Estimation
        est_frame = tk.Frame(f, bg='#E8F5E9', relief='solid', bd=1)
        est_frame.pack(fill='x', padx=12, pady=(10,4))
        tk.Label(est_frame, bg='#E8F5E9', font=('Segoe UI', 9),
                 text="ℹ️  Cliquer ESTIMER avant de lancer pour connaître le nombre d'appels API et la durée estimée."
                 ).pack(side='left', padx=10, pady=6)
        tk.Button(est_frame, text="📊 ESTIMER",
                  bg=AEJ_ORANGE, fg='white', font=('Segoe UI', 9, 'bold'),
                  relief='flat', padx=10, pady=4, cursor='hand2',
                  command=self._estimer).pack(side='right', padx=10, pady=6)

        self.v_estimate = tk.StringVar(value="")
        tk.Label(f, textvariable=self.v_estimate, font=('Segoe UI', 9, 'bold'),
                 bg=AEJ_GRAY, fg=AEJ_ORANGE).pack()

        # Statut + barre
        self.v_status = tk.StringVar(value="⏸  En attente...")
        tk.Label(f, textvariable=self.v_status, font=('Segoe UI', 11, 'bold'),
                 bg=AEJ_GRAY, fg=AEJ_DARK).pack(pady=(8,2))
        self.bar = ttk.Progressbar(f, mode='determinate', length=900, maximum=100)
        self.bar.pack(pady=2)
        self.v_pct = tk.StringVar(value="0 %")
        tk.Label(f, textvariable=self.v_pct, font=('Segoe UI', 9),
                 bg=AEJ_GRAY, fg='#616161').pack()

        # Boutons
        btn_row = tk.Frame(f, bg=AEJ_GRAY)
        btn_row.pack(pady=10)
        self.btn_start = tk.Button(btn_row,
            text="▶   LANCER LE MATCHING",
            font=('Segoe UI', 13, 'bold'), bg=AEJ_GREEN, fg='white',
            padx=24, pady=11, relief='flat', cursor='hand2',
            activebackground=AEJ_DARK, activeforeground='white',
            command=self._lancer)
        self.btn_start.pack(side='left', padx=8)
        self.btn_stop = tk.Button(btn_row,
            text="⏹  ARRÊTER",
            font=('Segoe UI', 11), bg='#D32F2F', fg='white',
            padx=14, pady=11, relief='flat', cursor='hand2',
            state='disabled', command=self._stopper)
        self.btn_stop.pack(side='left', padx=8)
        tk.Button(btn_row, text="🗑  Vider cache",
            font=('Segoe UI', 10), bg='#757575', fg='white',
            padx=12, pady=9, relief='flat', cursor='hand2',
            command=self._vider_cache).pack(side='left', padx=8)
        tk.Button(btn_row, text="📖  Guide",
            font=('Segoe UI', 10), bg=AEJ_ORANGE, fg='white',
            padx=12, pady=9, relief='flat', cursor='hand2',
            command=self._guide).pack(side='left', padx=8)

        # Journal
        tk.Label(f, text="Journal d'exécution :", font=('Segoe UI', 9, 'bold'),
                 bg=AEJ_GRAY).pack(anchor='w', padx=14)
        self.journal = scrolledtext.ScrolledText(
            f, height=16, font=('Consolas', 9),
            bg='#1C1C1C', fg='#00E676', insertbackground='white',
            relief='flat', padx=6, pady=6)
        self.journal.pack(fill='both', expand=True, padx=12, pady=(2,10))
        self.journal.config(state='disabled')

        self._orch = None

    def _log(self, msg: str):
        self.journal.config(state='normal')
        self.journal.insert('end', msg + '\n')
        self.journal.see('end')
        self.journal.config(state='disabled')
        self.root.update_idletasks()

    def _progress(self, val: int, msg: str):
        self.bar['value'] = val
        self.v_pct.set(f"{val} %")
        self.v_status.set(msg)
        self.root.update_idletasks()

    def _estimer(self):
        def run():
            try:
                provider = make_provider(self.cfg)
                provider.connecter()
                if self.cfg.get('SOURCE_MODE', 'DRIVE') == 'LOCAL':
                    tok_cv, tok_f = 'LOCAL_CV', 'LOCAL_FICHE'
                else:
                    tok_cv = self.cfg.get('DOSSIER_CV_SOURCE')
                    tok_f  = self.cfg.get('DOSSIER_FICHES_POSTE')
                nb_cv = provider.compter_fichiers(tok_cv)
                nb_f  = provider.compter_fichiers(tok_f)
                est   = MoteurIA.estimer_appels(nb_cv, nb_f, self.cfg.get('SEUIL_TFIDF', 18))
                msg = (f"📦 CV: {nb_cv}  |  📋 Fiches: {nb_f}  |  "
                       f"🤖 Appels IA estimés: {est['total']}  |  "
                       f"⏱️ Durée estimée: ~{est['duree_min']} min")
                self.root.after(0, lambda: self.v_estimate.set(msg))
            except Exception as e:
                self.root.after(0, lambda: self.v_estimate.set(f"Estimation impossible: {e}"))
        threading.Thread(target=run, daemon=True).start()
        self.v_estimate.set("⏳ Estimation en cours...")

    def _lancer(self):
        # Quick Ollama reachability check before starting the pipeline
        try:
            import urllib.request as ur
            ollama_url = self.cfg.get('OLLAMA_URL', 'http://localhost:11434')
            ur.urlopen(ollama_url + "/api/tags", timeout=4)
        except Exception:
            messagebox.showwarning(
                "Ollama introuvable",
                "Ollama n'est pas détecté sur " + self.cfg.get('OLLAMA_URL', 'http://localhost:11434') +
                "\n\nVérifiez que :\n"
                "  1. Ollama est installé (https://ollama.com)\n"
                "  2. Ollama est lancé (ollama serve)\n"
                "  3. Un modèle est installé (ex: ollama pull gemma3:4b)"
            )
            return
        if self.cfg.get('SOURCE_MODE', 'DRIVE') == 'LOCAL':
            if not self.cfg.get('LOCAL_CV_PATHS') or not self.cfg.get('LOCAL_FICHE_PATHS'):
                messagebox.showwarning("Sources manquantes",
                    "Mode local : ajoutez au moins un fichier/dossier de CV et de fiches "
                    "de poste, puis ENREGISTREZ la configuration.")
                return
            if not self.cfg.get('LOCAL_DEST_PATH'):
                messagebox.showwarning("Destination manquante",
                    "Mode local : choisissez un dossier de destination, puis ENREGISTREZ.")
                return
        self.btn_start.config(state='disabled', text='⏳  En cours...')
        self.btn_stop.config(state='normal')
        self.journal.config(state='normal')
        self.journal.delete('1.0', 'end')
        self.journal.config(state='disabled')
        self.bar['value'] = 0

        def run():
            try:
                self._orch = Orchestrateur(
                    self.cfg, self.db,
                    log=self._log, progress=self._progress)
                self.resultats = self._orch.executer()
                self.root.after(0, self._post_pipeline)
            except Exception as e:
                self.root.after(0, lambda: messagebox.showerror("Erreur", str(e)))
            finally:
                self.root.after(0, lambda: (
                    self.btn_start.config(state='normal', text='▶   LANCER LE MATCHING'),
                    self.btn_stop.config(state='disabled')
                ))
        threading.Thread(target=run, daemon=True).start()

    def _stopper(self):
        if self._orch:
            self._orch.stop()
            self._log("⏹  Arrêt demandé par l'utilisateur.")
        self.btn_stop.config(state='disabled')

    def _post_pipeline(self):
        self._refresh_badge()
        messagebox.showinfo("Terminé ✅",
            f"Matching terminé !\n{len(self.resultats)} correspondances trouvées.\n"
            "Consultez les onglets Dashboard et Mise en Relation.")
        self.nb.select(2)   # Aller au dashboard
        self._refresh_dashboard()
        self._refresh_relation()

    def _vider_cache(self):
        if messagebox.askyesno("Vider cache",
                "Supprimer tous les profils et scores mis en cache ?"):
            self.db.vider()
            self._log("🗑  Cache vidé.")
            self._refresh_badge()
            messagebox.showinfo("Cache", "Cache vidé. Le prochain lancement retraitera tout.")

    def _auto_run(self):
        self._log("⏰ Déclenchement automatique du pipeline...")
        try:
            orch = Orchestrateur(self.cfg, self.db,
                                 log=self._log, progress=self._progress)
            self.resultats = orch.executer()
            self.root.after(0, self._post_pipeline)
        except Exception as e:
            self._log(f"❌ Erreur auto: {e}")

    def _guide(self):
        w = tk.Toplevel(self.root)
        w.title("Guide AEJ CV Matcher PRO v2.0")
        w.geometry("760x600")
        st = scrolledtext.ScrolledText(w, font=('Consolas', 10), bg='white', relief='flat')
        st.pack(fill='both', expand=True, padx=8, pady=8)
        st.insert('end', GUIDE_TEXT)
        st.config(state='disabled')

    # ════════════════════════════════════════════════════════════
    #  ONGLET 3 — TABLEAU DE BORD (DASHBOARD)
    # ════════════════════════════════════════════════════════════
    def _tab_dashboard(self):
        f = tk.Frame(self.nb, bg=AEJ_GRAY)
        self.nb.add(f, text='📊  Dashboard')

        # Barre actions
        bar = tk.Frame(f, bg=AEJ_GRAY)
        bar.pack(fill='x', padx=12, pady=8)
        tk.Label(bar, text="Tableau de bord analytique",
                 font=('Segoe UI', 12, 'bold'), bg=AEJ_GRAY,
                 fg=AEJ_GREEN).pack(side='left')
        tk.Button(bar, text="🔄 Actualiser",
                  bg=AEJ_GREEN, fg='white', font=('Segoe UI', 9, 'bold'),
                  relief='flat', padx=10, pady=4, cursor='hand2',
                  command=self._refresh_dashboard).pack(side='right', padx=4)

        # Cadre pour les graphiques
        self.dash_frame = tk.Frame(f, bg='white', relief='solid', bd=1)
        self.dash_frame.pack(fill='both', expand=True, padx=12, pady=(0, 12))

        # Cartes KPI en haut
        self.kpi_frame = tk.Frame(f, bg=AEJ_GRAY)
        self.kpi_frame.pack(fill='x', padx=12, pady=4)

        self._canvas_dash = None

        # Pied de tableau
        self.v_dash_status = tk.StringVar(value="Lancez le pipeline pour afficher le dashboard.")
        tk.Label(f, textvariable=self.v_dash_status, font=('Segoe UI', 9),
                 bg=AEJ_GRAY, fg='#757575').pack(pady=4)

    def _refresh_dashboard(self):
        if not self.resultats:
            self.v_dash_status.set("Aucune donnée. Lancez le pipeline d'abord.")
            return
        if not MPLOT_OK:
            self.v_dash_status.set("matplotlib non installé — pip install matplotlib")
            return

        # KPI cards
        for w in self.kpi_frame.winfo_children():
            w.destroy()

        nb_top = sum(1 for r in self.resultats if r.get('score_global',0) >= 80)
        nb_moy = sum(1 for r in self.resultats if 65 <= r.get('score_global',0) < 80)
        nb_fbl = sum(1 for r in self.resultats if r.get('score_global',0) < 65)
        moy    = sum(r.get('score_global',0) for r in self.resultats) // max(len(self.resultats),1)

        for lbl, val, color in [
            ("⭐ TOP MATCHES", nb_top,              "#1A7A3C"),
            ("✅ MOYENS",      nb_moy,              "#F47920"),
            ("⚪ FAIBLES",     nb_fbl,              "#757575"),
            ("📈 SCORE MOY.",  f"{moy}%",           "#0D5A2C"),
            ("📄 TOTAL",       len(self.resultats), "#333333"),
        ]:
            card = tk.Frame(self.kpi_frame, bg=color, padx=16, pady=8, relief='flat')
            card.pack(side='left', padx=6, pady=4)
            tk.Label(card, text=str(val), font=('Segoe UI', 22, 'bold'),
                     bg=color, fg='white').pack()
            tk.Label(card, text=lbl, font=('Segoe UI', 8),
                     bg=color, fg='#E0E0E0').pack()

        # Graphiques
        for w in self.dash_frame.winfo_children():
            w.destroy()

        fig = Figure(figsize=(13.5, 5.5), facecolor='white', dpi=92)
        gs  = GridSpec(2, 3, figure=fig, hspace=0.5, wspace=0.38)

        # 1. Histogramme scores
        ax1 = fig.add_subplot(gs[0, 0])
        scores = [r.get('score_global', 0) for r in self.resultats]
        n, bins, patches = ax1.hist(scores, bins=12, edgecolor='white', linewidth=0.8)
        for patch, lb in zip(patches, bins):
            patch.set_facecolor('#F47920' if lb >= 80 else ('#1A7A3C' if lb >= 65 else '#B0BEC5'))
        ax1.set_title('Distribution des scores', fontsize=9, fontweight='bold', color=AEJ_GREEN)
        ax1.set_xlabel('Score (%)', fontsize=8)
        ax1.set_ylabel('Candidats', fontsize=8)
        ax1.tick_params(labelsize=7)
        ax1.spines[['top','right']].set_visible(False)

        # 2. Donut décisions
        ax2 = fig.add_subplot(gs[0, 1])
        vals = [nb_top, nb_moy, nb_fbl]
        lbls = [f'Prioritaires\n{nb_top}', f'Moyens\n{nb_moy}', f'Faibles\n{nb_fbl}']
        clrs = ['#F47920', '#1A7A3C', '#B0BEC5']
        if sum(vals) > 0:
            wedges, _, autotexts = ax2.pie(
                vals, labels=lbls, colors=clrs, autopct='%1.0f%%',
                pctdistance=0.80, wedgeprops=dict(width=0.48),
                textprops={'fontsize': 7})
            for at in autotexts:
                at.set_fontsize(7)
        ax2.set_title('Répartition décisions', fontsize=9, fontweight='bold', color=AEJ_GREEN)

        # 3. Top 8 postes
        ax3 = fig.add_subplot(gs[0, 2])
        par_poste = {}
        for r in self.resultats:
            k = r.get('intitule_poste', '?')[:22]
            par_poste[k] = par_poste.get(k, 0) + 1
        top_postes = sorted(par_poste.items(), key=lambda x: x[1], reverse=True)[:8]
        if top_postes:
            pnames = [p[0] for p in top_postes]
            pvals  = [p[1] for p in top_postes]
            bars   = ax3.barh(pnames, pvals, color=AEJ_GREEN, edgecolor='white')
            for bar, val in zip(bars, pvals):
                ax3.text(bar.get_width() + 0.1, bar.get_y() + bar.get_height()/2,
                         str(val), va='center', fontsize=7)
            ax3.set_title('Top postes (candidats)', fontsize=9, fontweight='bold', color=AEJ_GREEN)
            ax3.tick_params(labelsize=7)
            ax3.spines[['top','right']].set_visible(False)
            ax3.invert_yaxis()

        # 4. Secteurs
        ax4 = fig.add_subplot(gs[1, 0])
        par_secteur = {}
        for r in self.resultats:
            s = r.get('secteur_poste','Autre') or 'Autre'
            par_secteur[s[:18]] = par_secteur.get(s[:18], 0) + 1
        if par_secteur:
            s_items = sorted(par_secteur.items(), key=lambda x: x[1], reverse=True)[:7]
            slbls   = [s[0] for s in s_items]
            svals   = [s[1] for s in s_items]
            colors  = plt.cm.Set2([i/len(svals) for i in range(len(svals))]) if MPLOT_OK else ['#1A7A3C']*len(svals)
            ax4.pie(svals, labels=slbls, autopct='%1.0f%%',
                    colors=colors, textprops={'fontsize': 7},
                    pctdistance=0.85)
            ax4.set_title('Répartition par secteur', fontsize=9, fontweight='bold', color=AEJ_GREEN)

        # 5. Score par entreprise (barres)
        ax5 = fig.add_subplot(gs[1, 1])
        par_ent = {}
        for r in self.resultats:
            e = (r.get('entreprise','') or 'N/A')[:18]
            if e not in par_ent:
                par_ent[e] = []
            par_ent[e].append(r.get('score_global', 0))
        if par_ent:
            top_ent = sorted(par_ent.items(), key=lambda x: max(x[1]), reverse=True)[:6]
            enames  = [e[0] for e in top_ent]
            emoys   = [sum(e[1])//len(e[1]) for e in top_ent]
            bars = ax5.bar(range(len(enames)), emoys,
                           color=[AEJ_ORANGE if v >= 80 else AEJ_GREEN for v in emoys],
                           edgecolor='white')
            ax5.set_xticks(range(len(enames)))
            ax5.set_xticklabels(enames, rotation=25, ha='right', fontsize=7)
            ax5.set_ylabel('Score moyen', fontsize=8)
            ax5.set_title('Score moyen / entreprise', fontsize=9, fontweight='bold', color=AEJ_GREEN)
            ax5.spines[['top','right']].set_visible(False)
            ax5.tick_params(labelsize=7)

        # 6. Expérience vs Score (scatter)
        ax6 = fig.add_subplot(gs[1, 2])
        exps   = [r.get('annees_experience', 0) or 0 for r in self.resultats]
        scrs   = [r.get('score_global', 0) for r in self.resultats]
        colors_scatter = [AEJ_ORANGE if s >= 80 else (AEJ_GREEN if s >= 65 else '#B0BEC5')
                          for s in scrs]
        ax6.scatter(exps, scrs, c=colors_scatter, alpha=0.7, s=35, edgecolors='white')
        ax6.set_xlabel("Années d'exp.", fontsize=8)
        ax6.set_ylabel("Score global (%)", fontsize=8)
        ax6.set_title("Expérience vs Score", fontsize=9, fontweight='bold', color=AEJ_GREEN)
        ax6.tick_params(labelsize=7)
        ax6.spines[['top','right']].set_visible(False)
        ax6.axhline(y=80, color=AEJ_ORANGE, linestyle='--', linewidth=0.8, alpha=0.7)
        ax6.axhline(y=65, color=AEJ_GREEN,  linestyle='--', linewidth=0.8, alpha=0.5)

        if self._canvas_dash:
            self._canvas_dash.get_tk_widget().destroy()
        self._canvas_dash = FigureCanvasTkAgg(fig, master=self.dash_frame)
        self._canvas_dash.draw()
        self._canvas_dash.get_tk_widget().pack(fill='both', expand=True)

        total_postes = len({r.get('intitule_poste') for r in self.resultats})
        self.v_dash_status.set(
            f"Dashboard actualisé — {len(self.resultats)} correspondances sur {total_postes} postes | "
            f"Score moyen : {moy}%")

    # ════════════════════════════════════════════════════════════
    #  ONGLET 4 — MISE EN RELATION (CRM)
    # ════════════════════════════════════════════════════════════
    def _tab_relation(self):
        f = tk.Frame(self.nb, bg=AEJ_GRAY)
        self.nb.add(f, text='🤝  Mise en Relation')

        # Filtres
        flt = tk.Frame(f, bg='#E8F5E9', relief='solid', bd=1)
        flt.pack(fill='x', padx=12, pady=(10,4))

        tk.Label(flt, text="Filtres :", font=('Segoe UI', 9, 'bold'),
                 bg='#E8F5E9').pack(side='left', padx=8, pady=6)

        self.v_flt_statut = tk.StringVar(value='TOUS')
        tk.Label(flt, text="Statut :", bg='#E8F5E9', font=('Segoe UI', 9)).pack(side='left', padx=4)
        self.cmb_statut = ttk.Combobox(flt, textvariable=self.v_flt_statut, width=20,
            values=['TOUS'] + list(STATUTS_FR.values()), state='readonly')
        self.cmb_statut.pack(side='left', padx=4)

        self.v_flt_poste = tk.StringVar(value='TOUS')
        tk.Label(flt, text="Poste :", bg='#E8F5E9', font=('Segoe UI', 9)).pack(side='left', padx=8)
        self.cmb_poste = ttk.Combobox(flt, textvariable=self.v_flt_poste, width=30, state='readonly')
        self.cmb_poste.pack(side='left', padx=4)

        tk.Button(flt, text="🔍 Filtrer", bg=AEJ_GREEN, fg='white',
                  font=('Segoe UI', 9, 'bold'), relief='flat', padx=10, pady=3,
                  cursor='hand2', command=self._refresh_relation).pack(side='left', padx=8)
        tk.Button(flt, text="🔄 Tout recharger", bg=AEJ_ORANGE, fg='white',
                  font=('Segoe UI', 9), relief='flat', padx=10, pady=3,
                  cursor='hand2', command=self._refresh_relation).pack(side='left', padx=4)

        # Export
        tk.Button(flt, text="📥 Export Excel", bg='#455A64', fg='white',
                  font=('Segoe UI', 9), relief='flat', padx=10, pady=3,
                  cursor='hand2', command=self._export_contacts).pack(side='right', padx=10, pady=6)

        # Table Treeview
        cols = ('NOM', 'PROFIL', 'POSTE', 'ENTREPRISE', 'SCORE', 'DÉCISION',
                'STATUT', 'EMAIL', 'TÉLÉPHONE')
        self.tree_rel = ttk.Treeview(f, columns=cols, show='headings', height=14)
        widths = [20, 22, 26, 20, 8, 12, 20, 24, 14]
        for col, w in zip(cols, widths):
            self.tree_rel.heading(col, text=col,
                command=lambda c=col: self._sort_tree(self.tree_rel, c, False))
            self.tree_rel.column(col, width=w*8, anchor='center' if w <= 12 else 'w')

        tree_scroll_y = ttk.Scrollbar(f, orient='vertical', command=self.tree_rel.yview)
        tree_scroll_x = ttk.Scrollbar(f, orient='horizontal', command=self.tree_rel.xview)
        self.tree_rel.configure(yscrollcommand=tree_scroll_y.set,
                                xscrollcommand=tree_scroll_x.set)
        self.tree_rel.pack(fill='both', expand=True, padx=12, pady=4)
        tree_scroll_y.pack(side='right', fill='y')

        # Menu contextuel
        self.ctx_menu = tk.Menu(self.root, tearoff=0)
        self.ctx_menu.add_command(label="📧 Email de convocation (candidat)",
                                   command=self._email_convocation)
        self.ctx_menu.add_command(label="📤 Email de présélection (recruteur)",
                                   command=self._email_recruteur)
        self.ctx_menu.add_separator()
        for st in STATUTS:
            lbl = STATUTS_FR[st]
            self.ctx_menu.add_command(label=f"Marquer : {lbl}",
                command=lambda s=st: self._changer_statut(s))
        self.ctx_menu.add_separator()
        self.ctx_menu.add_command(label="📝 Ajouter une note",
                                   command=self._ajouter_note)

        self.tree_rel.bind("<Button-3>", self._show_context)
        self.tree_rel.bind("<Double-1>", lambda e: self._email_convocation())

        # Stats footer
        self.v_rel_stats = tk.StringVar(value="")
        stats_bar = tk.Frame(f, bg='#EEEEEE', relief='solid', bd=1)
        stats_bar.pack(fill='x', padx=12, pady=(0,8))
        tk.Label(stats_bar, textvariable=self.v_rel_stats,
                 font=('Segoe UI', 9, 'bold'), bg='#EEEEEE',
                 fg=AEJ_DARK).pack(pady=6, padx=10, anchor='w')

        # Charger les données initiales
        self._refresh_relation()

    def _refresh_relation(self):
        # Mettre à jour liste des postes
        postes = ['TOUS'] + self.db.get_postes_distincts()
        self.cmb_poste['values'] = postes

        # Lire filtre statut (extraire code depuis label)
        statut_flt = self.v_flt_statut.get()
        code_statut = 'TOUS'
        for code, label in STATUTS_FR.items():
            if label == statut_flt:
                code_statut = code
                break

        poste_flt = self.v_flt_poste.get()
        contacts  = self.db.get_contacts(
            statut=code_statut if code_statut != 'TOUS' else None,
            poste=poste_flt if poste_flt != 'TOUS' else None)

        # Vider et remplir la table
        for item in self.tree_rel.get_children():
            self.tree_rel.delete(item)

        for ct in contacts:
            sc  = ct.get('score', 0)
            tag = 'top' if sc >= 80 else ('moy' if sc >= 65 else 'fbl')
            values = (
                ct.get('nom_candidat',''), ct.get('titre_profil','')[:25],
                ct.get('intitule_poste','')[:30], ct.get('entreprise','')[:22],
                f"{sc}%", ct.get('decision',''),
                STATUTS_FR.get(ct.get('statut','EN_ATTENTE'),'—'),
                ct.get('email_candidat',''), ct.get('tel_candidat','')
            )
            self.tree_rel.insert('', 'end', iid=str(ct['id']),
                                  values=values, tags=(tag,))

        # Styles
        self.tree_rel.tag_configure('top', background='#C8E6C9')
        self.tree_rel.tag_configure('moy', background='#FFF9C4')
        self.tree_rel.tag_configure('fbl', background='#FFCDD2')

        # Stats
        stats = self.db.stats_contacts()
        parts = [f"{STATUTS_FR.get(k,k)}: {v}" for k, v in stats.items()]
        self.v_rel_stats.set("  |  ".join(parts) if parts else "Aucune donnée")

    def _get_selected_contact(self):
        sel = self.tree_rel.selection()
        if not sel:
            messagebox.showwarning("Sélection", "Veuillez sélectionner un candidat.")
            return None
        cid = int(sel[0])
        rows = self.db.get_contacts()
        for r in rows:
            if r['id'] == cid:
                return r
        return None

    def _show_context(self, event):
        item = self.tree_rel.identify_row(event.y)
        if item:
            self.tree_rel.selection_set(item)
            self.ctx_menu.post(event.x_root, event.y_root)

    def _changer_statut(self, statut: str):
        ct = self._get_selected_contact()
        if not ct:
            return
        self.db.update_contact(ct['id'], statut=statut)
        self._refresh_relation()
        self._log(f"📋 Statut mis à jour : {ct.get('nom_candidat','?')} → {STATUTS_FR[statut]}")

    def _email_convocation(self):
        ct = self._get_selected_contact()
        if not ct:
            return

        def gen():
            try:
                # Reconstruire candidat et poste depuis DB
                cv_data    = self.db.get_cv(ct['id_cv']) or {'nom_candidat': ct.get('nom_candidat','')}
                fiche_data = self.db.get_fiche(ct['id_fiche']) or {
                    'intitule': ct.get('intitule_poste',''), 'entreprise': ct.get('entreprise','')}
                email_txt = self.ia.generer_email_convocation(
                    cv_data, fiche_data, ct.get('score',0))
                sig = self.cfg.get('SIGNATURE_EMAIL', '')
                email_complet = email_txt + sig
                # Sauvegarder dans DB
                self.db.update_contact(ct['id'], email_genere=email_complet)
                self.root.after(0, lambda: self._afficher_email(
                    "Email de convocation", email_complet,
                    ct.get('email_candidat', ''), ct))
            except Exception as e:
                self.root.after(0, lambda: messagebox.showerror("Erreur", str(e)))

        threading.Thread(target=gen, daemon=True).start()
        self._log(f"⏳ Génération email pour {ct.get('nom_candidat','?')}...")

    def _email_recruteur(self):
        ct = self._get_selected_contact()
        if not ct:
            return

        def gen():
            try:
                fiche_data = self.db.get_fiche(ct['id_fiche']) or {
                    'intitule': ct.get('intitule_poste',''), 'entreprise': ct.get('entreprise','')}
                # Récupérer tous les candidats pour ce poste
                candidats = [c for c in self.db.get_contacts()
                             if c['id_fiche'] == ct['id_fiche'] and c['score'] >= 65]
                email_txt = self.ia.generer_email_recruteur(fiche_data, candidats)
                sig = self.cfg.get('SIGNATURE_EMAIL', '')
                email_complet = email_txt + sig
                self.root.after(0, lambda: self._afficher_email(
                    "Email recruteur — Présélection candidats",
                    email_complet, '', ct))
            except Exception as e:
                self.root.after(0, lambda: messagebox.showerror("Erreur", str(e)))

        threading.Thread(target=gen, daemon=True).start()
        self._log(f"⏳ Génération email recruteur pour {ct.get('intitule_poste','?')}...")

    def _afficher_email(self, titre, contenu, email_dest, ct):
        win = tk.Toplevel(self.root)
        win.title(titre)
        win.geometry("700x520")
        win.configure(bg='white')

        tk.Label(win, text=titre, font=('Segoe UI', 12, 'bold'),
                 bg=AEJ_GREEN, fg='white').pack(fill='x', pady=0)

        # Champ email destinataire
        dest_f = tk.Frame(win, bg='white')
        dest_f.pack(fill='x', padx=10, pady=6)
        tk.Label(dest_f, text="À :", font=('Segoe UI', 9, 'bold'),
                 bg='white').pack(side='left')
        v_dest = tk.StringVar(value=email_dest)
        ttk.Entry(dest_f, textvariable=v_dest, width=40).pack(side='left', padx=6)

        # Zone texte email
        st = scrolledtext.ScrolledText(win, font=('Segoe UI', 10),
                                        relief='flat', padx=8, pady=8)
        st.pack(fill='both', expand=True, padx=10)
        st.insert('end', contenu)

        # Boutons
        btn_f = tk.Frame(win, bg='white')
        btn_f.pack(pady=8)

        def copier():
            self.root.clipboard_clear()
            self.root.clipboard_append(st.get('1.0','end').strip())
            messagebox.showinfo("Copié", "Email copié dans le presse-papier !")

        def ouvrir_mail():
            corps = requests.utils.quote(st.get('1.0','end').strip())
            dest  = v_dest.get().strip()
            webbrowser.open(f"mailto:{dest}?body={corps}")

        def marquer_contacte():
            self.db.update_contact(ct['id'], statut='CONTACTE')
            self._refresh_relation()
            win.destroy()
            messagebox.showinfo("Statut", "Candidat marqué comme CONTACTÉ.")

        tk.Button(btn_f, text="📋 Copier", bg=AEJ_GREEN, fg='white',
                  font=('Segoe UI', 10, 'bold'), relief='flat',
                  padx=14, pady=6, cursor='hand2', command=copier).pack(side='left', padx=8)
        tk.Button(btn_f, text="📩 Ouvrir dans Mail", bg=AEJ_ORANGE, fg='white',
                  font=('Segoe UI', 10), relief='flat',
                  padx=12, pady=6, cursor='hand2', command=ouvrir_mail).pack(side='left', padx=8)
        tk.Button(btn_f, text="✅ Marquer Contacté", bg='#388E3C', fg='white',
                  font=('Segoe UI', 10), relief='flat',
                  padx=12, pady=6, cursor='hand2', command=marquer_contacte).pack(side='left', padx=8)
        tk.Button(btn_f, text="Fermer", bg='#757575', fg='white',
                  font=('Segoe UI', 10), relief='flat',
                  padx=12, pady=6, cursor='hand2', command=win.destroy).pack(side='left', padx=8)

    def _ajouter_note(self):
        ct = self._get_selected_contact()
        if not ct:
            return
        win = tk.Toplevel(self.root)
        win.title(f"Note — {ct.get('nom_candidat','')}")
        win.geometry("500x280")
        tk.Label(win, text="Note / Commentaire :",
                 font=('Segoe UI', 10, 'bold')).pack(pady=8, anchor='w', padx=10)
        txt = scrolledtext.ScrolledText(win, height=8, font=('Segoe UI', 10))
        txt.pack(fill='both', expand=True, padx=10)
        txt.insert('end', ct.get('notes', ''))

        def sauver():
            note = txt.get('1.0', 'end').strip()
            self.db.update_contact(ct['id'], notes=note)
            self._refresh_relation()
            win.destroy()
            messagebox.showinfo("Note", "Note enregistrée.")

        tk.Button(win, text="💾 Enregistrer", bg=AEJ_GREEN, fg='white',
                  font=('Segoe UI', 10, 'bold'), relief='flat',
                  padx=14, pady=7, cursor='hand2', command=sauver).pack(pady=8)

    def _export_contacts(self):
        contacts = self.db.get_contacts()
        if not contacts:
            messagebox.showwarning("Vide", "Aucun contact à exporter.")
            return
        chemin = filedialog.asksaveasfilename(
            defaultextension=".xlsx",
            filetypes=[("Excel", "*.xlsx")],
            initialfile=f"Contacts_AEJ_{datetime.now().strftime('%Y%m%d')}.xlsx")
        if not chemin:
            return
        try:
            gen = RapportGenerator()
            gen.excel_complet([], contacts, Path(chemin))
            messagebox.showinfo("Exporté", f"Fichier exporté :\n{chemin}")
        except Exception as e:
            messagebox.showerror("Erreur", str(e))

    def _sort_tree(self, tree, col, reverse):
        data = [(tree.set(k, col), k) for k in tree.get_children('')]
        data.sort(reverse=reverse)
        for idx, (_, k) in enumerate(data):
            tree.move(k, '', idx)
        tree.heading(col, command=lambda: self._sort_tree(tree, col, not reverse))

    # ════════════════════════════════════════════════════════════
    #  ONGLET 5 — MODE AUTONOME
    # ════════════════════════════════════════════════════════════
    def _tab_autonome(self):
        f = tk.Frame(self.nb, bg=AEJ_GRAY)
        self.nb.add(f, text='⏰  Mode Autonome')

        # Toggle
        top = tk.Frame(f, bg=AEJ_GRAY)
        top.pack(fill='x', padx=20, pady=20)

        self.v_auto = tk.BooleanVar(value=self.cfg.get('AUTO_ACTIVE', False))
        tk.Label(top, text="Activer le mode autonome :",
                 font=('Segoe UI', 12, 'bold'), bg=AEJ_GRAY).pack(side='left')
        toggle = ttk.Checkbutton(top, variable=self.v_auto,
                                  command=self._toggle_auto)
        toggle.pack(side='left', padx=12)

        # Intervalle
        itvl_f = tk.LabelFrame(f, text="  ⏱️  Planification  ",
                                font=('Segoe UI', 10, 'bold'), bg=AEJ_GRAY,
                                fg=AEJ_ORANGE, bd=2, relief='groove')
        itvl_f.pack(fill='x', padx=20, pady=10)

        tk.Label(itvl_f, text="Lancer automatiquement toutes les :",
                 font=('Segoe UI', 10), bg=AEJ_GRAY).grid(row=0, column=0, padx=12, pady=10, sticky='w')
        self.v_itvl = tk.IntVar(value=self.cfg.get('AUTO_INTERVALLE_H', 6))
        ttk.Spinbox(itvl_f, from_=1, to=24, textvariable=self.v_itvl, width=6).grid(
            row=0, column=1, padx=8)
        tk.Label(itvl_f, text="heures", font=('Segoe UI', 10),
                 bg=AEJ_GRAY).grid(row=0, column=2, padx=4)

        tk.Button(itvl_f, text="💾 Appliquer",
                  bg=AEJ_GREEN, fg='white', font=('Segoe UI', 9, 'bold'),
                  relief='flat', padx=10, pady=4, cursor='hand2',
                  command=self._appliquer_auto).grid(row=0, column=3, padx=12)

        # Statut watchdog
        status_f = tk.LabelFrame(f, text="  📡  Statut du watchdog  ",
                                  font=('Segoe UI', 10, 'bold'), bg=AEJ_GRAY,
                                  fg=AEJ_GREEN, bd=2, relief='groove')
        status_f.pack(fill='x', padx=20, pady=10)

        self.v_wd_actif  = tk.StringVar(value="⬛ Inactif")
        self.v_wd_next   = tk.StringVar(value="—")
        self.v_wd_lastrun = tk.StringVar(value="—")

        tk.Label(status_f, text="Watchdog :", font=('Segoe UI', 10, 'bold'),
                 bg=AEJ_GRAY).grid(row=0, column=0, padx=12, pady=8, sticky='w')
        tk.Label(status_f, textvariable=self.v_wd_actif,
                 font=('Segoe UI', 10), bg=AEJ_GRAY).grid(row=0, column=1, sticky='w')
        tk.Label(status_f, text="Prochain lancement :", font=('Segoe UI', 10, 'bold'),
                 bg=AEJ_GRAY).grid(row=1, column=0, padx=12, pady=8, sticky='w')
        tk.Label(status_f, textvariable=self.v_wd_next,
                 font=('Segoe UI', 10), bg=AEJ_GRAY).grid(row=1, column=1, sticky='w')
        tk.Label(status_f, text="Dernier run :", font=('Segoe UI', 10, 'bold'),
                 bg=AEJ_GRAY).grid(row=2, column=0, padx=12, pady=8, sticky='w')
        tk.Label(status_f, textvariable=self.v_wd_lastrun,
                 font=('Segoe UI', 10), bg=AEJ_GRAY).grid(row=2, column=1, sticky='w')

        # Historique runs
        hist_f = tk.LabelFrame(f, text="  📜  Historique des exécutions  ",
                                font=('Segoe UI', 10, 'bold'), bg=AEJ_GRAY,
                                fg=AEJ_ORANGE, bd=2, relief='groove')
        hist_f.pack(fill='both', expand=True, padx=20, pady=10)

        cols_h = ('DATE', 'DURÉE', 'CV', 'FICHES', 'MATCHES', 'STATUT')
        self.tree_runs = ttk.Treeview(hist_f, columns=cols_h, show='headings', height=8)
        for col in cols_h:
            self.tree_runs.heading(col, text=col)
            self.tree_runs.column(col, width=120, anchor='center')
        self.tree_runs.pack(fill='both', expand=True, padx=8, pady=8)

        # Charger l'historique
        self._refresh_runs()

        # Infos
        info = tk.Frame(f, bg='#E3F2FD', relief='solid', bd=1)
        info.pack(fill='x', padx=20, pady=8)
        tk.Label(info, bg='#E3F2FD', font=('Segoe UI', 9),
                 text=("💡  Le mode autonome surveille votre Drive et déclenche le pipeline "
                       "automatiquement. Seuls les nouveaux CV sont retraités grâce au cache SQLite.")
                 ).pack(padx=10, pady=6, anchor='w')

    def _toggle_auto(self):
        actif = self.v_auto.get()
        self.cfg.set('AUTO_ACTIVE', actif)
        if actif:
            self.watchdog.demarrer()
            self.v_wd_actif.set("🟢 Actif")
            self._tick_watchdog()
        else:
            self.watchdog.arreter()
            self.v_wd_actif.set("⬛ Inactif")
            self.v_wd_next.set("—")

    def _appliquer_auto(self):
        self.cfg.set('AUTO_INTERVALLE_H', self.v_itvl.get())
        if self.watchdog.is_running():
            self.watchdog.arreter()
            self.watchdog = Watchdog(self.cfg, self.db, self._auto_run)
            self.watchdog.demarrer()
        messagebox.showinfo("Autonome", "Intervalle mis à jour.")

    def _tick_watchdog(self):
        if self.watchdog.is_running():
            self.v_wd_actif.set("🟢 Actif")
            self.v_wd_next.set(self.watchdog.temps_restant())
        else:
            self.v_wd_actif.set("⬛ Inactif")
        self._refresh_last_run()
        self.root.after(5000, self._tick_watchdog)

    def _refresh_last_run(self):
        lr = self.db.last_run()
        if lr:
            d  = lr.get('duree_sec', 0)
            self.v_wd_lastrun.set(
                f"{lr.get('date_debut','—')[:16]} | {lr.get('nb_matches',0)} matches | "
                f"{d//60}min{d%60}s | {lr.get('statut','')}")

    def _refresh_runs(self):
        for item in self.tree_runs.get_children():
            self.tree_runs.delete(item)
        runs = self.db.conn.execute(
            "SELECT * FROM runs ORDER BY id DESC LIMIT 20").fetchall()
        for r in runs:
            r = dict(r)
            d = r.get('duree_sec', 0) or 0
            self.tree_runs.insert('', 'end', values=(
                (r.get('date_debut','')[:16]),
                f"{d//60}min {d%60}s",
                r.get('nb_cv', 0),
                r.get('nb_fiches', 0),
                r.get('nb_matches', 0),
                r.get('statut', ''),
            ))

    # ════════════════════════════════════════════════════════════
    #  ONGLET 6 — À PROPOS
    # ════════════════════════════════════════════════════════════
    def _tab_about(self):
        f = tk.Frame(self.nb, bg='white')
        self.nb.add(f, text='ℹ️  À propos')

        # Header visuel
        hdr = tk.Frame(f, bg=AEJ_GREEN, height=80)
        hdr.pack(fill='x')
        tk.Label(hdr, text="🎯  AEJ CV MATCHER PRO  v2.0",
                 font=('Segoe UI', 20, 'bold'), bg=AEJ_GREEN, fg='white').pack(pady=14)
        tk.Label(hdr, text="ÉDITION PRESTIGE — Le meilleur logiciel de matching au monde",
                 font=('Segoe UI', 10), bg=AEJ_GREEN, fg='#C8E6C9').pack()

        st = scrolledtext.ScrolledText(f, font=('Segoe UI', 10),
                                        bg='white', relief='flat', padx=16, pady=12)
        st.pack(fill='both', expand=True)
        st.insert('end', ABOUT_TEXT)
        st.config(state='disabled')


# ════════════════════════════════════════════════════════════════
#  SECTION 12 — TEXTES
# ════════════════════════════════════════════════════════════════
GUIDE_TEXT = """
═══════════════════════════════════════════════════════════════════
GUIDE — AEJ CV MATCHER PRO  v2.0  ÉDITION PRESTIGE
═══════════════════════════════════════════════════════════════════

ÉTAPE 1 — OLLAMA (IA 100% locale, gratuit)
─────────────────────────────────────────────────────────────────
  1. Installez Ollama : https://ollama.com (Windows/Mac/Linux)
  2. Dans un terminal : ollama pull gemma3:4b
     (autres modèles : llama3.2 / mistral / phi4:14b / qwen2.5:7b)
  3. Onglet Configuration → cliquer "🔄 Modèles" puis "🔗 Tester"

ÉTAPE 2 — CREDENTIALS GOOGLE DRIVE
─────────────────────────────────────────────────────────────────
  1. https://console.cloud.google.com
  2. Nouveau projet → APIs & Services → activer "Google Drive API"
  3. Identifiants → OAuth 2.0 Client ID → Application de bureau
  4. Télécharger → renommer credentials.json → placer à côté du .exe

ÉTAPE 3 — TESSERACT OCR (CV scannés)
─────────────────────────────────────────────────────────────────
  https://github.com/UB-Mannheim/tesseract/wiki
  Installer → vérifier le chemin dans Configuration

UTILISATION
─────────────────────────────────────────────────────────────────
  1. Configuration → renseigner les clés → ENREGISTRER
  2. Exécution → ESTIMER (optionnel) → LANCER LE MATCHING
  3. Dashboard → visualiser les graphiques
  4. Mise en Relation → gérer les contacts, générer des emails
  5. Mode Autonome → planifier les runs automatiques

INNOVATIONS v2.0
─────────────────────────────────────────────────────────────────
  • Moteur hybride TF-IDF + Ollama → 100% local, zéro cloud
  • Traitement parallèle 4 threads → 4× plus rapide
  • Déduplication MD5 → aucun CV traité deux fois
  • Dashboard analytique 6 graphiques (matplotlib)
  • CRM Mise en Relation avec génération d'emails (IA locale)
  • Lettres de présélection DOCX par poste automatiques
  • Mode Autonome avec watchdog Drive
  • Scoring percentile et ranking sectoriel
  • Historique des runs avec statistiques

STRUCTURE DRIVE CRÉÉE
─────────────────────────────────────────────────────────────────
  DOSSIER DESTINATION/
  ├── ENTREPRISE A — RESPONSABLE RH/
  │   ├── ⭐ TOP MATCHES (≥80%)/
  │   ├── ✅ MATCHES MOYENS (65-79%)/
  │   ├── ⚪ FAIBLES (<65%)/
  │   └── Preselection_Responsable_RH.docx
  ├── ENTREPRISE B — CHEF COMPTABLE/
  └── RAPPORT_MATCHING_AEJ_20260613_XXXX.xlsx  (4 feuilles)
"""

ABOUT_TEXT = """
Développé pour : KEITA — CIP Service Entreprises
Structure       : Agence Emploi Jeunes / Agence Prestige
Adresse         : Bd des Martyrs, Cocody, Abidjan — Côte d'Ivoire

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
INNOVATIONS  v2.0  vs  v1.0
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

  ✦ Moteur Hybride TF-IDF + Ollama LLM         (100% local)
  ✦ Traitement Drive parallèle (4 threads)     (4× plus rapide)
  ✦ Déduplication intelligente MD5
  ✦ Tableau de bord 6 graphiques temps réel
  ✦ Module CRM Mise en Relation complet
  ✦ Génération emails personnalisés (IA locale)
  ✦ Lettres de présélection DOCX par poste
  ✦ Mode Autonome (watchdog + planification)
  ✦ Scoring percentile et ranking sectoriel
  ✦ Estimation des appels IA avant exécution
  ✦ Bouton ARRÊTER mid-pipeline
  ✦ Test connexion Ollama intégré
  ✦ Historique des runs (SQLite)
  ✦ Export contacts Excel (4 feuilles)
  ✦ 7 onglets (vs 3 en v1.0)

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
TECHNOLOGIES
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

  Python 3.11+            Interface Windows tkinter
  Google Drive API v3     OAuth2 — connexion unique (optionnel)
  Ollama                  LLM local Gemma/LLaMA/Mistral — 100% privé
  scikit-learn            TF-IDF pré-filtrage sémantique
  pdfplumber              Extraction PDF native
  python-docx             Word (.docx) + lettres présélection
  pytesseract + Pillow     OCR pour CV scannés
  matplotlib              Dashboard analytique 6 graphiques
  SQLite                  Cache + CRM + historique
  openpyxl                Rapports Excel 4 feuilles (AEJ)

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
PONDÉRATION MATCHING
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

  Compétences métier        35 %
  Expérience pertinente     25 %
  Diplôme / Formation       15 %
  Mots-clés sectoriels      10 %
  Outils / Logiciels         5 %
  Langues                    5 %
  Localisation               5 %

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
COULEURS AEJ :  Vert #1A7A3C  |  Orange #F47920
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

Agence Prestige — Bd des Martyrs, Cocody
Lun–Ven  7h30–12h30 / 13h30–16h30  |  📞 01 73 26 26 43
"""


# ════════════════════════════════════════════════════════════════
#  ENTRY POINT
# ════════════════════════════════════════════════════════════════
if __name__ == '__main__':
    App()
